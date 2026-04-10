import os
import uuid
import tempfile
import subprocess
import base64
import requests
import json
from io import BytesIO
from datetime import datetime
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from docx_formatter import apply_guidelines
from google import genai
from google.genai import types
from pydantic import BaseModel, Field
import urllib3
from urllib.parse import quote
import time
from tenacity import retry, stop_after_attempt, wait_exponential, before_sleep_log
import logging
import anthropic

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("app")

urllib3.disable_warnings()

class SlideData(BaseModel):
    title: str = Field(
        description="Slide title: max ~10 words, single line when rendered (short headline)."
    )
    narrative: str = Field(
        description="1-2 short lines only (under ~220 chars); must be longer/more prominent than bullets.",
        default="",
    )
    punchline: str = Field(description="One takeaway line; unique per slide, placed at the bottom.")
    key_takeaway: str = Field(description="A single powerful sentence summarizing the strategic impact or core takeaway of the slide.", default="Strategic growth driver.")
    layout_type: str = Field(
        description="title_slide | section_divider | index_slide | title_and_content | two_column | diagram"
    )
    slide_archetype: str = Field(description="Must be one of: title, agenda, divider, standard, table, deep_dive, roadmap, options", default="standard")
    bullet_points: list[str] = Field(
        description="3-5 bullets; each bullet max ~120 chars — short lines for infographic rows beside icons."
    )
    bullet_icon_seeds: list[str] = Field(
        default_factory=list,
        description="One short English seed per bullet for DiceBear icons (same count as bullet_points when possible).",
    )
    table_data: list[list[str]] = Field(description="2D array of strings for table/matrix slides. First row is headers.", default=[])
    icon_keyword: str = Field(description="A single keyword for the AI-generated icon (DiceBear) that acts as the slide diagram; required on content slides.")
    keep_original_image: bool = Field(description="Set to true if the original image contains important visual data like a chart, diagram, or photo that should be kept on the slide.")


class SlideLayoutPlanItem(BaseModel):
    slide_index: int = Field(description="0-based index matching final deck order")
    layout_type: str = Field(description="title_slide | section_divider | index_slide | title_and_content | two_column | diagram")
    slide_archetype: str = Field(description="Archetype for this slide")
    purpose_one_line: str = Field(description="What this slide must accomplish in one line")
    visual_role: str = Field(description="How the AI icon will reinforce the message (one line)")


class PresentationLayoutPlan(BaseModel):
    deck_narrative: str = Field(description="2-5 sentences: story arc and flow of the entire deck")
    slides: list[SlideLayoutPlanItem] = Field(description="Exactly one entry per slide, in order")


class VisualQAItem(BaseModel):
    slide_index: int = Field(description="0-based slide index")
    severity: str = Field(description="low | medium | high")
    issue: str = Field(description="Specific layout, overlap, or content balance issue")


class VisualLayoutReviewResult(BaseModel):
    round_summary: str = Field(description="Brief summary of this review pass")
    issues: list[VisualQAItem] = Field(default_factory=list)
    suggested_fixes: list[str] = Field(description="Concrete fixes to apply to slide copy or structure", default_factory=list)


class ImageTextBlock(BaseModel):
    reading_order: int = Field(description="Order to read this block (1 = first)")
    text: str = Field(description="Verbatim or cleaned text from this region")
    region_hint: str = Field(
        description="Spatial hint, e.g. top_title, left_column, diagram_label, callout, footer"
    )


class ImageLayoutAnalysisPhase1(BaseModel):
    """Rich image-to-text: extract content and describe layout before rebuilding the slide."""
    full_text_reading_order: str = Field(
        description="All readable text in natural reading order, one block per line or paragraph"
    )
    extracted_text_blocks: list[ImageTextBlock] = Field(
        description="Every distinct text region with position hints"
    )
    layout_description: str = Field(
        description="Detailed description of how the slide looks: zones, alignment, columns, spacing, hierarchy"
    )
    diagram_structure: str = Field(
        description="Faithful description of diagrams: boxes, arrows, flows, charts, connectors, grouping"
    )
    color_and_style_notes: str = Field(
        default="",
        description="Notable colors, emphasis, icons or photos visible",
    )
    visual_motifs_for_icons: list[str] = Field(
        description="5-12 short English keywords for AI-generated icons matching motifs (shapes, metaphors) seen in the image",
        default_factory=list,
    )


class TextBoxNorm(BaseModel):
    """Editable text region; coordinates normalized 0-1 over the slide (origin top-left)."""
    reading_order: int = Field(default=0)
    left: float = Field(ge=0.0, le=1.0)
    top: float = Field(ge=0.0, le=1.0)
    width: float = Field(ge=0.0, le=1.0)
    height: float = Field(ge=0.0, le=1.0)
    text: str
    font_emphasis: str = Field(
        default="normal",
        description="normal | bold | small_caption",
    )


class ImageToPptReconstruction(BaseModel):
    """Structured spec to build an editable slide that mirrors the source image."""
    title: str = Field(default="")
    narrative: str = Field(default="")
    punchline: str = Field(default="")
    bullet_points: list[str] = Field(default_factory=list)
    layout_type: str = Field(
        default="title_and_content",
        description="title_and_content | two_column | diagram",
    )
    icon_keyword: str = Field(
        default="presentation",
        description="Primary DiceBear icon seed aligned to the main visual metaphor",
    )
    extra_icon_keywords: list[str] = Field(
        default_factory=list,
        description="Additional icon seeds for secondary motifs (placed in a row)",
    )
    text_boxes: list[TextBoxNorm] = Field(
        default_factory=list,
        description="All text as separate editable boxes approximating positions from the image",
    )
    place_original_image_as_reference: bool = Field(
        default=True,
        description="If true, embed the source image in reference_image_box for diagram/chart fidelity",
    )
    reference_image_box: TextBoxNorm | None = Field(
        default=None,
        description="Where to place the source image; if null and place_original_image_as_reference, engine uses default right panel",
    )


class PresentationData(BaseModel):
    slides: list[SlideData] = Field(description="List of generated slides")

class SectionData(BaseModel):
    heading: str = Field(description="Section heading")
    content: str = Field(description="Multiple paragraphs of text for this section")
    bullet_points: list[str] = Field(description="Optional bullet points for this section")

class DocumentData(BaseModel):
    title: str = Field(description="Document title")
    sections: list[SectionData] = Field(description="Document sections")

THEMES = {
    "dark corporate": {
        "bg": (28, 30, 38),          # Deep slate grey background
        "accent": (0, 161, 241),    # Vivid blue accent
        "text": (245, 245, 245),    # Off-white text
        "subtext": (170, 175, 185)  # Dimmed text
    },
    "modern light": {
        "bg": (250, 252, 255),      # Crisp very-light blue/white
        "accent": (230, 57, 70),    # Strong red accent
        "text": (33, 37, 41),       # Almost black text
        "subtext": (108, 117, 125)  # Grey subtext
    },
    "pastel": {
        "bg": (245, 241, 237),      # Soft cream
        "accent": (148, 201, 169),  # Mint green
        "text": (73, 80, 87),       # Soft dark grey text
        "subtext": (140, 145, 150)  # Lighter grey
    },
    "blue accent": {
        "bg": (255, 255, 255),      # Pure white
        "accent": (0, 80, 158),     # Navy blue
        "text": (10, 25, 47),       # Very dark blue text
        "subtext": (80, 90, 110)    # Mid blue-grey
    },
    # Default light canvas: cool neutral background, indigo ribbons, strong title/punchline contrast
    "presentation light": {
        "bg": (248, 250, 252),
        "accent": (79, 70, 229),
        "title": (15, 23, 42),
        "text": (51, 65, 85),
        "subtext": (100, 116, 139),
        "punchline": (4, 120, 87),
    },
    # Dark canvas: slate background, warm titles, green punchlines
    "presentation dark": {
        "bg": (30, 41, 59),
        "accent": (251, 191, 36),
        "title": (251, 191, 36),
        "text": (226, 232, 240),
        "subtext": (148, 163, 184),
        "punchline": (34, 197, 94),
    },
}

# Canonical default for API/MCP; resolves to `presentation light` colors in `_get_theme_colors`.
DEFAULT_LAYOUT_THEME = "Studio Light"


def _is_dark_studio_theme(theme_str: str) -> bool:
    """Dark default pair (Studio Dark / Presentation Dark); legacy alias: voiceqa."""
    t = (theme_str or "").lower()
    return (
        "presentation dark" in t
        or "studio dark" in t
        or "voiceqa" in t
    )


def _wants_split_visual_layout(theme_str: str) -> bool:
    """Two-column text + visual column; optional — not the default row+hero infographic."""
    t = (theme_str or "").lower()
    if any(
        k in t
        for k in (
            "split-visual",
            "split visual",
            "split layout",
            "split-panel",
            "split panel",
            "two-panel",
            "two panel",
            "two-column visual",
            "two column visual",
        )
    ):
        return True
    return False


def _get_theme_colors(theme_str: str):
    t = (theme_str or "").lower().strip()
    if _is_dark_studio_theme(theme_str):
        return THEMES["presentation dark"]
    if "studio light" in t or "presentation light" in t:
        return THEMES["presentation light"]
    if "modern light" in t:
        return THEMES["modern light"]
    if "dark corporate" in t:
        return THEMES["dark corporate"]
    if "pastel" in t:
        return THEMES["pastel"]
    if "blue" in t:
        return THEMES["blue accent"]
    if "dark" in t:
        return THEMES["dark corporate"]
    return THEMES["presentation light"]

def _create_themed_presentation(theme_str: str):
    """
    Creates a new Presentation() and injects a robust, stylish layout template 
    into the master slide so every slide inherits a consistent, beautiful design.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    colors = _get_theme_colors(theme_str)
    bg_color = RGBColor(*colors["bg"])
    accent_color = RGBColor(*colors["accent"])
    
    # Apply to Slide Master so it inherits everywhere
    master = prs.slide_master
    
    # 1. Set solid background
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = bg_color
    
    return prs, colors

def _apply_theme_ribbons(slide, prs, colors):
    """Adds the thematic ribbons to an individual slide since python-pptx doesn't support adding shapes to masters."""
    from pptx.enum.shapes import MSO_SHAPE
    accent_color = RGBColor(*colors["accent"])
    
    top_ribbon = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        0, 0, prs.slide_width, Inches(0.15)
    )
    top_ribbon.fill.solid()
    top_ribbon.fill.fore_color.rgb = accent_color
    top_ribbon.line.fill.background()
    
    bottom_ribbon = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        0, prs.slide_height - Inches(0.3), prs.slide_width, Inches(0.3)
    )
    bottom_ribbon.fill.solid()
    bottom_ribbon.fill.fore_color.rgb = accent_color
    bottom_ribbon.line.fill.background()

def _apply_aptos_narrow(shape, font_color=None):
    if not hasattr(shape, 'text_frame'):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Aptos Narrow'
            if font_color:
                run.font.color.rgb = font_color


# Prompt-generated deck: hierarchy title > narrative > punchline > bullets (never larger body than narrative)
DECK_TITLE_PT = 22
DECK_NARRATIVE_PT = 15
DECK_PUNCHLINE_PT = 13
DECK_BULLET_PT = 11
DECK_TITLE_MAX_CHARS = 72


def _truncate_one_line_title(text: str, max_chars: int = DECK_TITLE_MAX_CHARS) -> str:
    t = (text or "").replace("\n", " ").strip()
    if len(t) <= max_chars:
        return t
    cut = t[: max_chars - 1].rsplit(" ", 1)[0]
    if len(cut) < max_chars // 2:
        cut = t[: max_chars - 1]
    return cut.rstrip() + "…"


def _dicebear_icon_url(seed: str, bg_hex: str = "e8e8e8") -> str:
    s = quote((seed or "visual").strip()[:80], safe="")
    bg = "".join(c for c in bg_hex.lower() if c in "0123456789abcdef")[:6] or "e8e8e8"
    return f"https://api.dicebear.com/9.x/icons/png?seed={s}&backgroundColor={bg}&size=256"


def _download_dicebear_icon(seed: str, dest_path: str, bg_hex: str = "e8e8e8") -> bool:
    try:
        r = requests.get(_dicebear_icon_url(seed, bg_hex=bg_hex), verify=False, timeout=20)
        if r.status_code != 200 or len(r.content) < 80:
            logger.warning("DiceBear HTTP %s len=%s", r.status_code, len(r.content) if r.content else 0)
            return False
        with open(dest_path, "wb") as f:
            f.write(r.content)
        return True
    except Exception as e:
        logger.warning("DiceBear download failed: %s", e)
    return False


def _style_slide_title_shape(
    shape,
    raw_text: str,
    text_color: RGBColor,
    *,
    truncate: bool = True,
    font_pt: int | None = None,
) -> None:
    shape.text = _truncate_one_line_title(raw_text) if truncate else (raw_text or "").strip()
    tf = shape.text_frame
    tf.word_wrap = False
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    pt = font_pt if font_pt is not None else DECK_TITLE_PT
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(pt)
        p.font.bold = True
        p.font.color.rgb = text_color
    _apply_aptos_narrow(shape, font_color=text_color)


def _deck_render_profile(theme_colors: dict, layout_theme: str) -> dict:
    """Default: infographic rows with per-bullet + hero DiceBear icons. Optional split layout when theme asks for it."""
    profile = {
        "split_visual": False,
        "title_pt": DECK_TITLE_PT,
        "narrative_pt": DECK_NARRATIVE_PT,
        "punchline_pt": DECK_PUNCHLINE_PT,
        "bullet_pt": DECK_BULLET_PT,
        "punchline_center": False,
        "punchline_bold": False,
        "dicebear_bg": "e8e8e8",
    }
    if _is_dark_studio_theme(layout_theme):
        profile.update(
            {
                "title_pt": 28,
                "narrative_pt": 16,
                "punchline_pt": 14,
                "bullet_pt": 12,
                "punchline_center": True,
                "punchline_bold": True,
                "dicebear_bg": "475569",
            }
        )
    if _wants_split_visual_layout(layout_theme):
        profile["split_visual"] = True
    return profile


def _title_color_from_theme(theme_colors: dict, text_color: RGBColor) -> RGBColor:
    if "title" in theme_colors:
        return RGBColor(*theme_colors["title"])
    return text_color


def _punchline_color_from_theme(theme_colors: dict) -> RGBColor:
    if "punchline" in theme_colors:
        return RGBColor(*theme_colors["punchline"])
    return RGBColor(*theme_colors["subtext"])


def _add_punchline_box(
    slide,
    punchline: str,
    theme_colors: dict,
    profile: dict,
) -> None:
    px = slide.shapes.add_textbox(
        Inches(0.5), SLIDE_HEIGHT - Inches(0.85), SLIDE_WIDTH - Inches(1.0), Inches(0.48)
    )
    pr = px.text_frame.paragraphs[0]
    pr.text = (punchline or "")[:420]
    pr.font.italic = True
    pr.font.bold = profile.get("punchline_bold", False)
    pr.font.size = Pt(profile["punchline_pt"])
    pr.font.color.rgb = _punchline_color_from_theme(theme_colors)
    pr.alignment = PP_ALIGN.CENTER if profile.get("punchline_center") else PP_ALIGN.LEFT
    _apply_aptos_narrow(px, font_color=pr.font.color.rgb)


def _add_strict_content_slide_rows(
    slide,
    s_data: dict,
    run_dir: str,
    slide_idx: int,
    theme_colors: dict,
    text_color: RGBColor,
    profile: dict,
    title_color: RGBColor,
    db_bg: str,
) -> None:
    """Infographic rows: icon + bullet text; hero icon top-right."""
    narrative_width = SLIDE_WIDTH - Inches(1.0) - Inches(1.38)
    title_raw = s_data.get("title", f"Slide {slide_idx + 1}")
    tit = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), narrative_width, Inches(0.62))
    _style_slide_title_shape(
        tit, title_raw, title_color, truncate=True, font_pt=profile["title_pt"]
    )

    nx = slide.shapes.add_textbox(Inches(0.5), Inches(0.92), narrative_width, Inches(0.55))
    ntf = nx.text_frame
    ntf.word_wrap = True
    np = ntf.paragraphs[0]
    np.text = (s_data.get("narrative") or "")[:520]
    np.font.size = Pt(profile["narrative_pt"])
    np.font.color.rgb = text_color
    np.alignment = PP_ALIGN.LEFT
    _apply_aptos_narrow(nx, font_color=text_color)

    bullets = [b.strip() for b in (s_data.get("bullet_points") or []) if isinstance(b, str) and b.strip()][:5]
    seeds_in = [str(x).strip() for x in (s_data.get("bullet_icon_seeds") or []) if str(x).strip()]
    main_seed = (s_data.get("icon_keyword") or "insight").strip() or "insight"

    row_top_in = 1.58
    usable = 6.52 - row_top_in
    nrows = max(1, len(bullets))
    row_h_in = min(0.92, usable / min(nrows, 5)) if bullets else 0.9

    for j, bullet in enumerate(bullets[:5]):
        bt = row_top_in + j * row_h_in
        seed = seeds_in[j] if j < len(seeds_in) else f"{main_seed}-{j + 1}"
        ip = os.path.join(run_dir, f"row_icon_{slide_idx}_{j}.png")
        if _download_dicebear_icon(seed, ip, bg_hex=db_bg):
            slide.shapes.add_picture(ip, Inches(0.48), Inches(bt), Inches(0.74), Inches(0.74))
        tx = slide.shapes.add_textbox(
            Inches(1.34),
            Inches(bt + 0.02),
            narrative_width - Inches(0.82),
            Inches(max(0.36, row_h_in - 0.05)),
        )
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet[:380]
        p.font.size = Pt(profile["bullet_pt"])
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.LEFT
        _apply_aptos_narrow(tx, font_color=text_color)

    _add_punchline_box(slide, s_data.get("punchline") or "", theme_colors, profile)

    hero = os.path.join(run_dir, f"hero_icon_{slide_idx}.png")
    if _download_dicebear_icon(main_seed, hero, bg_hex=db_bg):
        slide.shapes.add_picture(hero, Inches(10.95), Inches(0.28), Inches(1.38), Inches(1.38))


def _add_strict_content_slide_split(
    slide,
    s_data: dict,
    run_dir: str,
    slide_idx: int,
    theme_colors: dict,
    text_color: RGBColor,
    profile: dict,
    title_color: RGBColor,
    db_bg: str,
) -> None:
    """Optional split: narrative + bullets in left column; icon grid + hero on the right."""
    left_w = Inches(6.15)
    title_raw = s_data.get("title", f"Slide {slide_idx + 1}")
    tit = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), left_w, Inches(0.68))
    _style_slide_title_shape(
        tit, title_raw, title_color, truncate=True, font_pt=profile["title_pt"]
    )

    bullets = [b.strip() for b in (s_data.get("bullet_points") or []) if isinstance(b, str) and b.strip()][:5]
    seeds_in = [str(x).strip() for x in (s_data.get("bullet_icon_seeds") or []) if str(x).strip()]
    main_seed = (s_data.get("icon_keyword") or "theme").strip() or "theme"

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), left_w, Inches(5.15))
    tf = body.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = (s_data.get("narrative") or "")[:520]
    p0.font.size = Pt(profile["narrative_pt"])
    p0.font.color.rgb = text_color
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet[:320]}"
        p.level = 0
        p.font.size = Pt(profile["bullet_pt"])
        p.font.color.rgb = text_color
    _apply_aptos_narrow(body, font_color=text_color)

    _add_punchline_box(slide, s_data.get("punchline") or "", theme_colors, profile)

    # Right visual column (hero + small grid)
    hx = Inches(6.85)
    hero_path = os.path.join(run_dir, f"split_hero_{slide_idx}.png")
    if _download_dicebear_icon(main_seed, hero_path, bg_hex=db_bg):
        slide.shapes.add_picture(hero_path, hx, Inches(0.38), Inches(2.55), Inches(2.55))

    grid_seeds = []
    for i in range(4):
        if i < len(seeds_in) and seeds_in[i]:
            grid_seeds.append(seeds_in[i])
        else:
            grid_seeds.append(f"{main_seed}-{i + 1}")
    positions = [(6.9, 3.15), (9.35, 3.15), (6.9, 4.85), (9.35, 4.85)]
    for gi, (sx, sy) in enumerate(positions):
        if gi >= len(grid_seeds):
            break
        gp = os.path.join(run_dir, f"split_grid_{slide_idx}_{gi}.png")
        if _download_dicebear_icon(grid_seeds[gi], gp, bg_hex=db_bg):
            slide.shapes.add_picture(gp, Inches(sx), Inches(sy), Inches(1.15), Inches(1.15))


def _add_strict_content_slide_infographic(
    slide,
    s_data: dict,
    run_dir: str,
    slide_idx: int,
    theme_colors: dict,
    text_color: RGBColor,
    theme_low: str,
    layout_theme: str,
) -> None:
    """Blank-layout slide: optional split columns, or default row infographic with AI icons."""
    profile = _deck_render_profile(theme_colors, layout_theme)
    title_color = _title_color_from_theme(theme_colors, text_color)
    db_bg = profile["dicebear_bg"]
    if profile["split_visual"]:
        _add_strict_content_slide_split(
            slide,
            s_data,
            run_dir,
            slide_idx,
            theme_colors,
            text_color,
            profile,
            title_color,
            db_bg,
        )
    else:
        _add_strict_content_slide_rows(
            slide,
            s_data,
            run_dir,
            slide_idx,
            theme_colors,
            text_color,
            profile,
            title_color,
            db_bg,
        )


@retry(
    stop=stop_after_attempt(5),
    wait=wait_exponential(multiplier=1, min=4, max=60),
    before_sleep=before_sleep_log(logger, logging.WARNING),
    reraise=True
)
def _call_genai_with_retry(client, pil_img, prompt_text, schema: type[BaseModel] = SlideData):
    try:
        return client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[pil_img, prompt_text],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=schema,
                temperature=0.2
            ),
        )
    except Exception as e:
        if "429" in str(e) or "quota" in str(e).lower() or "rate" in str(e).lower():
            logger.warning(f"GenAI rate limit hit (429/Quota): {e}")
        else:
            logger.warning(f"GenAI API error: {e}")
        raise e

@retry(
    stop=stop_after_attempt(5),
    wait=wait_exponential(multiplier=1, min=4, max=60),
    before_sleep=before_sleep_log(logger, logging.WARNING),
    reraise=True
)
def _call_anthropic_with_retry(client, b64_img, prompt_text, schema: type[BaseModel] = SlideData):
    try:
        schema_hint = json.dumps(schema.model_json_schema(), indent=2)
        response = client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=8192,
            temperature=0.2,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": "image/png",
                                "data": b64_img
                            }
                        },
                        {
                            "type": "text",
                            "text": prompt_text + "\n\nRespond ONLY with a valid JSON object matching this schema:\n" + schema_hint
                        }
                    ]
                }
            ]
        )
        return response.content[0].text
    except Exception as e:
        if "429" in str(e) or "quota" in str(e).lower() or "rate" in str(e).lower():
            logger.warning(f"Anthropic rate limit hit (429/Quota): {e}")
        else:
            logger.warning(f"Anthropic API error: {e}")
        raise e

@retry(
    stop=stop_after_attempt(5),
    wait=wait_exponential(multiplier=1, min=4, max=60),
    before_sleep=before_sleep_log(logger, logging.WARNING),
    reraise=True
)
def _call_genai_text_with_retry(client, prompt_text, schema):
    try:
        return client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[prompt_text],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=schema,
                temperature=0.4
            ),
        )
    except Exception as e:
        if "429" in str(e) or "quota" in str(e).lower() or "rate" in str(e).lower():
            logger.warning(f"GenAI text rate limit hit: {e}")
        raise e

@retry(
    stop=stop_after_attempt(5),
    wait=wait_exponential(multiplier=1, min=4, max=60),
    before_sleep=before_sleep_log(logger, logging.WARNING),
    reraise=True
)
def _call_anthropic_text_with_retry(client, prompt_text, schema):
    try:
        response = client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=8192,
            temperature=0.4,
            messages=[
                {
                    "role": "user",
                    "content": prompt_text + "\n\nRespond ONLY with a valid JSON object matching this schema:\n" + json.dumps(schema.model_json_schema(), indent=2)
                }
            ]
        )
        return response.content[0].text
    except Exception as e:
        if "429" in str(e) or "quota" in str(e).lower() or "rate" in str(e).lower():
            logger.warning(f"Anthropic text rate limit hit: {e}")
        raise e

def _fit_image_to_slide(slide, img_path, slide_width, slide_height, margin):
    img = Image.open(img_path)
    img_width, img_height = img.size
    page_aspect = img_width / img_height
    slide_aspect = (slide_width - 2 * margin) / (slide_height - 2 * margin)
    
    if page_aspect > slide_aspect:
        width = slide_width - 2 * margin
        height = width / page_aspect
    else:
        height = slide_height - 2 * margin
        width = height * page_aspect
        
    left = (slide_width - width) / 2 + margin
    top = (slide_height - height) / 2 + margin
    
    slide.shapes.add_picture(img_path, left, top, width, height)

def _trigger_webhook(webhook_url: str, payload: dict):
    """
    Push the generated artifact info to the specified webhook URL.
    """
    if not webhook_url:
        return
    try:
        # Avoid hanging the generation response by using a short timeout for webhook
        requests.post(webhook_url, json=payload, timeout=10)
    except Exception as e:
        print(f"Failed to trigger webhook at {webhook_url}: {e}")

def format_document(doc_source: str, is_url: bool = True, webhook_url: str = None) -> dict:
    """
    Downloads/Reads a DOCX file, applies corporate guidelines, and returns the formatted DOCX URL.
    """
    stats["requests_received"] += 1
    stats["last_request_time"] = datetime.now().isoformat()
    
    execution_id = str(uuid.uuid4())
    run_dir = os.path.join(OUTPUT_DIR, execution_id)
    os.makedirs(run_dir, exist_ok=True)
    
    input_filename = "input.docx"
    input_path = os.path.join(run_dir, input_filename)
    output_filename = "formatted_document.docx"
    output_path = os.path.join(run_dir, output_filename)
    
    try:
        if is_url:
            if doc_source.startswith(('http://', 'https://')):
                response = requests.get(doc_source, verify=False)
                response.raise_for_status()
                with open(input_path, 'wb') as f:
                    f.write(response.content)
            else:
                if ";base64," in doc_source:
                    _, b64_data = doc_source.split(";base64,")
                    with open(input_path, 'wb') as f:
                        f.write(base64.b64decode(b64_data))
                else:
                    import shutil
                    shutil.copy2(doc_source, input_path)
        else:
            with open(input_path, 'wb') as f:
                f.write(base64.b64decode(doc_source))
                
        # Apply formatting guidelines
        apply_guidelines(input_path, output_path)
        
        file_url = _get_file_url(execution_id, output_filename)
        _record_success(file_url, output_filename)
        return {
            "success": True,
            "message": "Document formatted successfully.",
            "file_url": file_url,
            "execution_id": execution_id,
            "filename": output_filename
        }
        
    except Exception as e:
        stats["failed_creations"] += 1
        return {
            "success": False,
            "message": f"Error formatting document: {str(e)}"
        }

def _send_progress(webhook_url, message, status="in_progress"):
    if not webhook_url:
        return
    try:
        requests.post(webhook_url, json={"status": status, "message": message}, verify=False, timeout=5)
    except:
        pass

def process_pdf_to_artifacts(
    pdf_source: str, 
    is_url: bool = True, 
    instructions: str = "", 
    layout_theme: str = "", 
    visual_iconography: str = "", 
    slide_content_rules: str = "",
    target_format: str = "pptx",
    webhook_url: str = None,
    api_key: str = ""
) -> dict:
    """
    Converts a PDF into a PPTX or DOCX, incorporating custom guidelines.
    """
    stats["requests_received"] += 1
    stats["last_request_time"] = datetime.now().isoformat()
    
    execution_id = str(uuid.uuid4())
    run_dir = os.path.join(OUTPUT_DIR, execution_id)
    os.makedirs(run_dir, exist_ok=True)
    
    input_filename = "input.pdf"
    input_path = os.path.join(run_dir, input_filename)
    
    try:
        # Load PDF
        if is_url:
            if pdf_source.startswith(('http://', 'https://')):
                response = requests.get(pdf_source, verify=False)
                response.raise_for_status()
                with open(input_path, 'wb') as f:
                    f.write(response.content)
            else:
                if ";base64," in pdf_source:
                    _, b64_data = pdf_source.split(";base64,")
                    with open(input_path, 'wb') as f:
                        f.write(base64.b64decode(b64_data))
                else:
                    import shutil
                    shutil.copy2(pdf_source, input_path)
        else:
            with open(input_path, 'wb') as f:
                f.write(base64.b64decode(pdf_source))
                
        doc = fitz.open(input_path)
        try:
            if target_format.lower() == "pptx":
                output_filename = "converted_presentation.pptx"
                output_path = os.path.join(run_dir, output_filename)
                
                prs, theme_colors = _create_themed_presentation(layout_theme)
                
                # Add an instructions slide to pass metadata or indicate rules
                if instructions or layout_theme or visual_iconography or slide_content_rules:
                    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
                    slide.shapes.title.text = "Generated PPTX Guidelines applied"
                    tf = slide.placeholders[1].text_frame
                    tf.text = "The following guidelines were requested for this presentation:\n"
                    if layout_theme: tf.add_paragraph().text = f"- Theme: {layout_theme}"
                    if visual_iconography: tf.add_paragraph().text = f"- Iconography: {visual_iconography}"
                    if slide_content_rules: tf.add_paragraph().text = f"- Content Rules: {slide_content_rules}"
                    if instructions: tf.add_paragraph().text = f"- Instructions: {instructions}"
                    
                blank_layout = prs.slide_layouts[6]
                
                # Determine which AI client to use
                has_genai = False
                has_anthropic = False
                client = None
                
                # Check for explicit key or environment key
                use_anthropic = False
                if api_key.startswith("sk-ant") or (not api_key and os.environ.get("ANTHROPIC_API_KEY")):
                    use_anthropic = True
                    
                if use_anthropic:
                    try:
                        k = api_key if api_key else os.environ.get("ANTHROPIC_API_KEY")
                        proxy_url = os.environ.get("GCP_PROXY_FOR_CLAUD")
                        if proxy_url:
                            client = anthropic.Anthropic(api_key=k, base_url=proxy_url, max_retries=0)
                        else:
                            client = anthropic.Anthropic(api_key=k, max_retries=0)
                        has_anthropic = True
                    except:
                        pass
                else:
                    try:
                        if api_key:
                            client = genai.Client(api_key=api_key)
                        elif os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY"):
                            if not os.environ.get("GCP_PROXY_FOR_CLAUD"):
                                client = genai.Client()
                        if client:
                            has_genai = True
                    except:
                        pass
                        
                ai_rate_limit_fallback_count = 0

                for page_num in range(len(doc)):
                    _send_progress(webhook_url, f"Processing page {page_num + 1} of {len(doc)}...")
                    page = doc[page_num]
                    # Generate a high-res image for AI to analyze
                    mat = fitz.Matrix(3.0, 3.0)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    img_path = os.path.join(run_dir, f"page_{page_num}.png")
                    pix.save(img_path)
                    
                    if has_genai or has_anthropic:
                        try:
                            prompt_text = f"""Analyze this presentation slide image thoroughly. 
                            1. Extract all text content and structure it logically.
                            2. Suggest a new layout type (choose strictly from: title_and_content, two_column, diagram).
                            3. Provide an overarching punchline.
                            4. Provide a single keyword for a visual icon that represents the core idea.
                            5. Determine if the original image contains complex diagrams, charts, or essential visual data that MUST be kept on the slide (set keep_original_image to true if so).
                            
                            User Instructions: {instructions}
                            Layout Theme: {layout_theme}
                            Content Rules: {slide_content_rules}"""
                            
                            try:
                                if has_anthropic:
                                    with open(img_path, "rb") as image_file:
                                        b64_img = base64.b64encode(image_file.read()).decode('utf-8')
                                    prompt_text += f"\n\nJSON SCHEMA:\n{json.dumps(SlideData.model_json_schema(), indent=2)}"
                                    raw_text = _call_anthropic_with_retry(client, b64_img, prompt_text)
                                else:
                                    pil_img = Image.open(img_path)
                                    response = _call_genai_with_retry(client, pil_img, prompt_text)
                                    raw_text = response.text.strip()
                            except Exception as api_err:
                                print(f"AI API rate limit or other error after retries: {api_err}")
                                raise api_err

                            try:
                                # Clean response string just in case it has markdown code block formatting
                                if raw_text.startswith("```json"):
                                    raw_text = raw_text[7:]
                                elif raw_text.startswith("```"):
                                    raw_text = raw_text[3:]
                                if raw_text.endswith("```"):
                                    raw_text = raw_text[:-3]
                                slide_data = json.loads(raw_text.strip())
                            except Exception as parse_e:
                                print(f"JSON Parse error for page {page_num}: {parse_e}\nRaw output was: {raw_text}")
                                raise parse_e
                            
                            # Build editable slide
                            l_type = slide_data.get('layout_type', 'title_and_content')
                            if l_type == 'diagram':
                                slide_layout = prs.slide_layouts[6] # Blank
                                slide = prs.slides.add_slide(slide_layout)
                                
                                # Apply theme ribbons
                                _apply_theme_ribbons(slide, prs, theme_colors)
                                
                                # We keep the original image for diagrams
                                _fit_image_to_slide(slide, img_path, SLIDE_WIDTH, SLIDE_HEIGHT, MARGIN)
                                
                                # Add a title text box over the image
                                left = Inches(0.5)
                                top = Inches(0.2)
                                width = Inches(12.0)
                                height = Inches(0.5)
                                txBox = slide.shapes.add_textbox(left, top, width, height)
                                tf = txBox.text_frame
                                p = tf.add_paragraph()
                                p.text = slide_data.get('title', f"Slide {page_num + 1}")
                                p.font.bold = True
                                p.font.size = Pt(28)
                                _apply_aptos_narrow(txBox)
                            else:
                                if l_type == 'two_column':
                                    slide_layout = prs.slide_layouts[3] # Two Content
                                else:
                                    slide_layout = prs.slide_layouts[1] # Title and Content
                                    
                                slide = prs.slides.add_slide(slide_layout)
                                
                                # Set Title
                                title_shape = slide.shapes.title
                                title_shape.left = Inches(0.5)
                                title_shape.top = Inches(0.25)
                                title_shape.width = SLIDE_WIDTH - Inches(1.0)
                                title_shape.height = Inches(0.8)
                                _apply_aptos_narrow(title_shape, font_color=RGBColor(*theme_colors["text"]))
                                
                                # Set Narrative
                                left = Inches(0.5)
                                top = Inches(0.95)
                                width = SLIDE_WIDTH - Inches(1.0)
                                height = Inches(0.5)
                                txBox = slide.shapes.add_textbox(left, top, width, height)
                                tf = txBox.text_frame
                                p = tf.add_paragraph()
                                p.text = slide_data.get('narrative', '')
                                p.font.size = Pt(16)
                                p.font.color.rgb = RGBColor(*theme_colors["text"])
                                _apply_aptos_narrow(txBox)
                                
                                # Set Punchline at bottom
                                left = Inches(0.5)
                                top = SLIDE_HEIGHT - Inches(0.8)
                                width = SLIDE_WIDTH - Inches(1.0)
                                height = Inches(0.4)
                                txBox_punch = slide.shapes.add_textbox(left, top, width, height)
                                tf_punch = txBox_punch.text_frame
                                p = tf_punch.add_paragraph()
                                p.text = slide_data.get('punchline', '')
                                p.font.italic = True
                                p.font.size = Pt(14)
                                p.font.color.rgb = RGBColor(*theme_colors["subtext"])
                                _apply_aptos_narrow(txBox_punch)
                                
                                # Set Bullets
                                body_shape = slide.placeholders[1]
                                body_shape.left = Inches(0.5)
                                body_shape.top = Inches(1.6)
                                body_shape.width = SLIDE_WIDTH - Inches(1.0)
                                body_shape.height = Inches(4.8)
                                tf = body_shape.text_frame
                                tf.word_wrap = True
                                tf.text = "" # clear default
                                for bullet in slide_data.get('bullet_points', []):
                                    p = tf.add_paragraph()
                                    p.text = bullet
                                    p.level = 0
                                _apply_aptos_narrow(body_shape, font_color=RGBColor(*theme_colors["text"]))
                                
                                # Add AI Generated Icon
                                icon_keyword = slide_data.get('icon_keyword', 'presentation')
                                icon_path = os.path.join(run_dir, f"icon_{page_num}.png")
                                if _download_dicebear_icon(icon_keyword, icon_path):
                                    slide.shapes.add_picture(icon_path, Inches(11.5), Inches(0.5), Inches(1), Inches(1))
                                    
                                # If two column, or if AI indicated we should keep the image, put it in the right placeholder
                                keep_image = slide_data.get('keep_original_image', False)
                                if (l_type == 'two_column' or keep_image) and len(slide.placeholders) > 2:
                                    # Adjust left body shape to be half width
                                    body_shape.width = (SLIDE_WIDTH / 2) - Inches(0.75)
                                    
                                    right_body_shape = slide.placeholders[2]
                                    right_body_shape.left = (SLIDE_WIDTH / 2) + Inches(0.25)
                                    right_body_shape.top = Inches(1.8)
                                    right_body_shape.width = (SLIDE_WIDTH / 2) - Inches(0.75)
                                    right_body_shape.height = Inches(5.0)
                                    tf_right = right_body_shape.text_frame
                                    tf_right.word_wrap = True
                                    tf_right.text = "Original Context"
                                    _apply_aptos_narrow(right_body_shape, font_color=RGBColor(*theme_colors["text"]))
                                    
                                    # Insert the original image overlapping the placeholder slightly
                                    slide.shapes.add_picture(
                                        img_path, 
                                        right_body_shape.left, 
                                        right_body_shape.top + Inches(0.5), 
                                        right_body_shape.width, 
                                        right_body_shape.height - Inches(0.5)
                                    )
                                elif keep_image:
                                    # Just add it to the bottom right
                                    slide.shapes.add_picture(img_path, Inches(8.0), Inches(4.0), Inches(4.5), Inches(3.0))
                        except Exception as e:
                            print(f"GenAI failed for page {page_num}: {e}")
                            ai_rate_limit_fallback_count += 1
                            
                            # Fallback if AI completely fails
                            slide = prs.slides.add_slide(blank_layout)
                            _fit_image_to_slide(slide, img_path, SLIDE_WIDTH, SLIDE_HEIGHT, MARGIN)
                    else:
                        slide = prs.slides.add_slide(blank_layout)
                        _fit_image_to_slide(slide, img_path, SLIDE_WIDTH, SLIDE_HEIGHT, MARGIN)
                        
                    os.remove(img_path)
                    
                prs.save(output_path)
            else:
                # DOCX
                output_filename = "converted_document.docx"
                output_path = os.path.join(run_dir, output_filename)
                
                docx_doc = DocxDocument()
                docx_doc.add_heading('Generated Document from PDF', 0)
                
                if instructions or layout_theme or visual_iconography or slide_content_rules:
                    docx_doc.add_heading('Generation Guidelines', level=1)
                    if layout_theme: docx_doc.add_paragraph(f"Theme: {layout_theme}", style='List Bullet')
                    if visual_iconography: docx_doc.add_paragraph(f"Iconography: {visual_iconography}", style='List Bullet')
                    if slide_content_rules: docx_doc.add_paragraph(f"Content Rules: {slide_content_rules}", style='List Bullet')
                    if instructions: docx_doc.add_paragraph(f"Instructions: {instructions}", style='List Bullet')
                    
                docx_doc.add_heading('Extracted Content', level=1)
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    text = page.get_text("text")
                    if text.strip():
                        docx_doc.add_paragraph(text)
                
                docx_doc.save(output_path)
                
                # Apply corporate guidelines to the generated docx
                formatted_output_filename = "final_formatted_document.docx"
                formatted_output_path = os.path.join(run_dir, formatted_output_filename)
                apply_guidelines(output_path, formatted_output_path)
                output_filename = formatted_output_filename
        finally:
            doc.close()
            
        file_url = _get_file_url(execution_id, output_filename)
        stats["successful_creations"] += 1
        _add_to_history(execution_id, output_filename, file_url, "process_pdf")
        
        msg = f"Successfully generated {target_format.upper()} from PDF."
        if target_format.lower() == "pptx" and not (has_genai or has_anthropic):
            msg += " Note: No valid API key found. Fell back to generating static image slides."
        elif target_format.lower() == "pptx" and ai_rate_limit_fallback_count > 0:
            msg += f" Note: {ai_rate_limit_fallback_count} slides fell back to original images due to AI API rate limits or errors. Retries were attempted."

        response_payload = {
            "success": True,
            "message": msg,
            "file_url": file_url,
            "download_path": f"/downloads/{execution_id}/{output_filename}",
            "execution_id": execution_id,
            "filename": output_filename
        }
        _trigger_webhook(webhook_url, response_payload)
        return response_payload
        
    except Exception as e:
        stats["failed_creations"] += 1
        error_payload = {
            "success": False,
            "message": f"Error converting PDF: {str(e)}"
        }
        _trigger_webhook(webhook_url, error_payload)
        return error_payload
    """
    Downloads/Reads a DOCX file, applies corporate guidelines, and returns the formatted DOCX URL.
    """
    stats["requests_received"] += 1
    stats["last_request_time"] = datetime.now().isoformat()
    
    execution_id = str(uuid.uuid4())
    run_dir = os.path.join(OUTPUT_DIR, execution_id)
    os.makedirs(run_dir, exist_ok=True)
    
    input_filename = "input.docx"
    input_path = os.path.join(run_dir, input_filename)
    output_filename = "formatted_document.docx"
    output_path = os.path.join(run_dir, output_filename)
    
    try:
        if is_url:
            if doc_source.startswith(('http://', 'https://')):
                response = requests.get(doc_source, verify=False)
                response.raise_for_status()
                with open(input_path, 'wb') as f:
                    f.write(response.content)
            else:
                if ";base64," in doc_source:
                    _, b64_data = doc_source.split(";base64,")
                    with open(input_path, 'wb') as f:
                        f.write(base64.b64decode(b64_data))
                else:
                    import shutil
                    shutil.copy2(doc_source, input_path)
        else:
            with open(input_path, 'wb') as f:
                f.write(base64.b64decode(doc_source))
                
        # Apply formatting guidelines
        apply_guidelines(input_path, output_path)
        
        file_url = _get_file_url(execution_id, output_filename)
        
        stats["successful_creations"] += 1
        _add_to_history(execution_id, output_filename, file_url, "format_docx")
        
        response_payload = {
            "success": True,
            "message": "Document formatted successfully.",
            "file_url": file_url,
            "download_path": f"/downloads/{execution_id}/{output_filename}",
            "execution_id": execution_id,
            "filename": output_filename
        }
        _trigger_webhook(webhook_url, response_payload)
        return response_payload
        
    except Exception as e:
        stats["failed_creations"] += 1
        error_payload = {
            "success": False,
            "message": f"Error formatting document: {str(e)}"
        }
        _trigger_webhook(webhook_url, error_payload)
        return error_payload

# Global stats
stats = {
    "requests_received": 0,
    "successful_creations": 0,
    "failed_creations": 0,
    "last_request_time": None,
    "last_success_file_url": None,
    "last_success_filename": None
}

# Global history for last X generated artifacts
MAX_HISTORY_ITEMS = int(os.environ.get("MAX_HISTORY_ITEMS", "10"))
generation_history = []

def _add_to_history(execution_id: str, filename: str, file_url: str, artifact_type: str):
    """Add a successful generation to the history list."""
    history_item = {
        "execution_id": execution_id,
        "filename": filename,
        "file_url": file_url,
        "type": artifact_type,
        "timestamp": datetime.now().isoformat()
    }
    generation_history.insert(0, history_item) # Add to front (newest first)
    
    # Keep only the last X items
    while len(generation_history) > MAX_HISTORY_ITEMS:
        generation_history.pop()

OUTPUT_DIR = os.environ.get("PPTX_OUTPUT_DIR", tempfile.gettempdir())
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.25)

# Layout types / archetypes that do NOT require full content-slide anatomy (title + narrative + icon + punchline)
_NON_CONTENT_LAYOUT_TYPES = frozenset({"title_slide", "section_divider", "index_slide"})
_NON_CONTENT_ARCHETYPES = frozenset({"title", "divider", "agenda"})


def _is_strict_content_slide(s: dict) -> bool:
    """True when the slide must have title, narrative, AI icon (diagram), and punchline."""
    lt = (s.get("layout_type") or "").lower().strip()
    arch = (s.get("slide_archetype") or "").lower().strip()
    if lt in _NON_CONTENT_LAYOUT_TYPES:
        return False
    if arch in _NON_CONTENT_ARCHETYPES:
        return False
    return True


def _validate_strict_content_slides(slides: list) -> list[str]:
    """Return human-readable validation errors for content-slide contract."""
    errors = []
    for i, s in enumerate(slides or []):
        if not isinstance(s, dict):
            errors.append(f"Slide {i}: invalid slide object.")
            continue
        if not _is_strict_content_slide(s):
            continue
        title = (s.get("title") or "").strip()
        narrative = (s.get("narrative") or "").strip()
        punchline = (s.get("punchline") or "").strip()
        icon_kw = (s.get("icon_keyword") or "").strip()
        arch = (s.get("slide_archetype") or "").lower().strip()
        bullets = s.get("bullet_points") or []
        table_data = s.get("table_data") or []

        if not title:
            errors.append(f"Slide {i}: content slide requires a non-empty title.")
        if not narrative:
            errors.append(f"Slide {i}: content slide requires a non-empty narrative.")
        if not punchline:
            errors.append(f"Slide {i}: content slide requires a non-empty punchline.")
        if not icon_kw:
            errors.append(f"Slide {i}: content slide requires icon_keyword for the AI-generated diagram/icon.")

        if len(title) > 78:
            errors.append(
                f"Slide {i}: title is too long for a single-line layout; shorten to about 72 characters."
            )

        bullet_seeds = s.get("bullet_icon_seeds") or []
        if arch in ("table", "roadmap", "options") and table_data and len(table_data) >= 2:
            pass
        else:
            non_empty_bullets = [b for b in bullets if isinstance(b, str) and b.strip()]
            if len(non_empty_bullets) < 2:
                errors.append(
                    f"Slide {i}: content slide needs at least two substantive bullet points "
                    f"(or a populated table_data for table/roadmap/options slides)."
                )
            else:
                n_seeds = len([x for x in bullet_seeds if str(x).strip()])
                if n_seeds < len(non_empty_bullets):
                    errors.append(
                        f"Slide {i}: bullet_icon_seeds must include one non-empty seed per bullet "
                        f"(infographic row icons); need {len(non_empty_bullets)}, have {n_seeds}."
                    )
    return errors


def _layout_regions_text_for_qa() -> str:
    """Fixed layout coordinates (inches) used by the renderer — for AI visual-layout QA."""
    return (
        "Slide canvas: 13.333 x 7.5 inches (16:9). "
        "Regions for standard content slides: "
        "title box ~left 0.5 top 0.25 width ~12.83 height 0.8; "
        "narrative textbox ~left 0.5 top 0.95 width ~11.18 (right margin reserved for 1\" icon + gap); "
        "content slides use infographic rows: small icon ~0.74in per bullet left column, bullet text to the right; "
        "large hero icon ~top-right 1.38in; punchline at bottom. "
        "Flag overlap or unreadable density if bullets are too long."
    )


def _llm_json_structured(client, use_anthropic: bool, prompt: str, schema: type[BaseModel]) -> dict:
    """Run structured JSON generation; returns a dict (Gemini or Anthropic)."""
    if use_anthropic:
        raw = _call_anthropic_text_with_retry(client, prompt, schema)
        if isinstance(raw, str):
            return json.loads(raw.strip())
        return json.loads(str(raw))
    response = _call_genai_text_with_retry(client, prompt, schema)
    text = getattr(response, "text", None)
    if not text:
        raise Exception("LLM returned empty response for structured JSON")
    return json.loads(text.strip())


def _repair_presentation_slides(
    client,
    use_anthropic: bool,
    slides: list,
    validation_errors: list[str],
    layout_plan: dict | None,
    visual_notes: list[str] | None,
) -> list:
    """Ask the LLM to fix slide JSON while preserving deck order and slide count."""
    plan_txt = json.dumps(layout_plan, indent=2) if layout_plan else "{}"
    ve = "\n".join(validation_errors) if validation_errors else "(none — if visual notes exist, apply those.)"
    vn = "\n".join(visual_notes or []) or "(none)"
    prompt = f"""You are fixing a presentation JSON. The slides array must keep the SAME length and order (same number of slides).

VALIDATION ERRORS TO FIX:
{ve}

VISUAL / LAYOUT NOTES FROM QA:
{vn}

APPROVED LAYOUT PLAN (must still be respected):
{plan_txt}

CURRENT SLIDES JSON:
{json.dumps(slides, indent=2)}

Rules for content slides (not title_slide, section_divider, index_slide, and not archetypes title/agenda/divider):
- Every such slide MUST have: short single-line title (~72 chars max), concise narrative (~220 chars), punchline, icon_keyword, bullet_points (3-5 short bullets ~120 chars each), and bullet_icon_seeds with the SAME count as bullet_points (one icon seed per bullet for infographic rows).
- Do not remove slides. Do not add slides. Fix fields only.

Return ONLY JSON matching the PresentationData schema (top key \"slides\")."""
    data = _llm_json_structured(client, use_anthropic, prompt, PresentationData)
    return data.get("slides") or slides


def _run_visual_layout_review(
    client,
    use_anthropic: bool,
    slides: list,
    layout_plan: dict | None,
    round_idx: int,
    previous: VisualLayoutReviewResult | None,
) -> VisualLayoutReviewResult:
    """One AI pass reviewing logical layout, crowding, and overlap risk from structured specs (no raster image)."""
    prev_txt = ""
    if previous and previous.issues:
        prev_txt = f"\nPrevious round issues to verify or resolve:\n{json.dumps([i.model_dump() for i in previous.issues], indent=2)}\n"
    prompt = f"""You are a senior presentation visual QA reviewer (round {round_idx} of 2).
Review the deck for layout quality: overlap risk between title, narrative, body, punchline, and the top-right AI icon diagram.
{_layout_regions_text_for_qa()}

Deck narrative / plan:
{json.dumps(layout_plan, indent=2) if layout_plan else "N/A"}

Slides data:
{json.dumps(slides, indent=2)}
{prev_txt}
Identify concrete issues (severity high if overlap or unreadable crowding likely). Suggest specific fixes to copy or structure.
Be strict on content slides: they must remain complete (title, narrative, bullets/table, icon keyword, punchline).
Return JSON matching the schema."""
    data = _llm_json_structured(client, use_anthropic, prompt, VisualLayoutReviewResult)
    return VisualLayoutReviewResult.model_validate(data)


def _get_file_url(execution_id: str, filename: str) -> str:
    file_path = os.path.abspath(os.path.join(OUTPUT_DIR, execution_id, filename))
    base_url = os.environ.get("BASE_URL", "")
    
    if base_url == "file://":
        return f"file:///{file_path.replace(chr(92), '/')}"
    elif base_url:
        prefix = base_url.rstrip('/')
        return f"{prefix}/downloads/{execution_id}/{filename}"
    else:
        # Fallback absolute path if no BASE_URL is set
        return f"http://localhost:8000/downloads/{execution_id}/{filename}"

def generate_presentation(python_code: str, webhook_url: str = None) -> dict:
    """
    Executes Python code to generate a PPTX file.
    Returns a dict with 'success', 'message', and optionally 'file_url' or 'execution_id'.
    """
    stats["requests_received"] += 1
    stats["last_request_time"] = datetime.now().isoformat()
    
    execution_id = str(uuid.uuid4())
    run_dir = os.path.join(OUTPUT_DIR, execution_id)
    os.makedirs(run_dir, exist_ok=True)
    
    script_path = os.path.join(run_dir, "script.py")
    
    with open(script_path, "w", encoding="utf-8") as f:
        f.write(python_code)
        
    try:
        result = subprocess.run(
            ["python", "script.py"],
            cwd=run_dir, 
            capture_output=True, 
            text=True, 
            timeout=60
        )
        
        if result.returncode != 0:
            stats["failed_creations"] += 1
            return {
                "success": False,
                "message": f"Error executing code:\n{result.stderr}\n\nStdout:\n{result.stdout}"
            }
            
        pptx_files = [f for f in os.listdir(run_dir) if f.endswith(".pptx")]
        
        if not pptx_files:
            stats["failed_creations"] += 1
            return {
                "success": False,
                "message": "Execution succeeded but no .pptx file was found. Ensure your code calls presentation.save('output.pptx').\n\nStdout:\n" + result.stdout
            }
            
        file_url = _get_file_url(execution_id, pptx_files[0])
            
        stats["successful_creations"] += 1
        _add_to_history(execution_id, pptx_files[0], file_url, "generate_pptx")
        
        response_payload = {
            "success": True,
            "message": "Presentation generated successfully.",
            "file_url": file_url,
            "download_path": f"/downloads/{execution_id}/{pptx_files[0]}",
            "execution_id": execution_id,
            "filename": pptx_files[0]
        }
        _trigger_webhook(webhook_url, response_payload)
        return response_payload
        
    except subprocess.TimeoutExpired:
        stats["failed_creations"] += 1
        error_payload = {
            "success": False,
            "message": "Error: Python code execution timed out after 60 seconds."
        }
        _trigger_webhook(webhook_url, error_payload)
        return error_payload
    except Exception as e:
        stats["failed_creations"] += 1
        error_payload = {
            "success": False,
            "message": f"Error: {str(e)}"
        }
        _trigger_webhook(webhook_url, error_payload)
        return error_payload

def _norm_rect_to_inches(tb: TextBoxNorm):
    """Map normalized 0-1 rectangle to slide inches (13.333 x 7.5)."""
    l = max(0.0, min(1.0, tb.left))
    t = max(0.0, min(1.0, tb.top))
    w = max(0.02, min(1.0, tb.width))
    h = max(0.02, min(1.0, tb.height))
    return (
        Inches(13.333 * l),
        Inches(7.5 * t),
        Inches(max(0.4, 13.333 * w)),
        Inches(max(0.35, 7.5 * h)),
    )


def _build_presentation_from_image_reconstruction(
    run_dir: str,
    source_img_path: str,
    recon: ImageToPptReconstruction,
    layout_theme: str,
) -> Presentation:
    """Build one editable slide: optional reference image panel + text boxes + AI icons."""
    prs, theme_colors = _create_themed_presentation(layout_theme)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg_color = RGBColor(*theme_colors["bg"])
    text_color = RGBColor(*theme_colors["text"])
    theme_low = layout_theme.lower()
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    _apply_theme_ribbons(slide, prs, theme_colors)

    # Reference thumbnail of source (diagram fidelity)
    if recon.place_original_image_as_reference and os.path.isfile(source_img_path):
        ref = recon.reference_image_box
        if ref is None:
            ref = TextBoxNorm(
                reading_order=-1,
                left=0.52,
                top=0.08,
                width=0.46,
                height=0.84,
                text="",
            )
        rl, rt, rw, rh = _norm_rect_to_inches(ref)
        try:
            slide.shapes.add_picture(source_img_path, rl, rt, rw, rh)
        except Exception as e:
            logger.warning(f"Could not embed reference image: {e}")

    # AI-generated icons (DiceBear), aligned to motifs from the image analysis
    icon_seeds = []
    if (recon.icon_keyword or "").strip():
        icon_seeds.append(recon.icon_keyword.strip())
    for k in recon.extra_icon_keywords or []:
        k = (k or "").strip()
        if k and k not in icon_seeds:
            icon_seeds.append(k)
    icon_seeds = icon_seeds[:8]

    ref_on_right = bool(recon.place_original_image_as_reference)
    for idx, seed in enumerate(icon_seeds):
        icon_url = f"https://api.dicebear.com/9.x/icons/png?seed={quote(str(seed), safe='')}&backgroundColor=ffffff"
        try:
            icon_resp = requests.get(icon_url, verify=False, timeout=15)
            if icon_resp.status_code == 200:
                icon_path = os.path.join(run_dir, f"img2ppt_icon_{idx}.png")
                with open(icon_path, "wb") as f:
                    f.write(icon_resp.content)
                if ref_on_right:
                    lx = Inches(0.45 + idx * 1.05)
                    ly = Inches(0.42)
                else:
                    lx = Inches(11.2 - min(idx, 3) * 1.05)
                    ly = Inches(0.45)
                slide.shapes.add_picture(icon_path, lx, ly, Inches(0.95), Inches(0.95))
        except Exception as e:
            logger.warning(f"Icon fetch failed for {seed}: {e}")

    # Editable text boxes (approximate positions from phase-2 reconstruction)
    boxes = sorted(recon.text_boxes or [], key=lambda b: b.reading_order)
    if not boxes and (recon.title or recon.narrative or recon.bullet_points):
        # Minimal fallback layout
        y = 0.12
        if recon.title:
            boxes.append(
                TextBoxNorm(
                    reading_order=1,
                    left=0.05,
                    top=y,
                    width=0.9 if not ref_on_right else 0.45,
                    height=0.1,
                    text=recon.title,
                    font_emphasis="bold",
                )
            )
            y += 0.12
        if recon.narrative:
            boxes.append(
                TextBoxNorm(
                    reading_order=2,
                    left=0.05,
                    top=y,
                    width=0.9 if not ref_on_right else 0.45,
                    height=0.08,
                    text=recon.narrative,
                )
            )
            y += 0.1
        for bi, bullet in enumerate(recon.bullet_points or []):
            boxes.append(
                TextBoxNorm(
                    reading_order=10 + bi,
                    left=0.05,
                    top=min(0.72, y + bi * 0.07),
                    width=0.9 if not ref_on_right else 0.45,
                    height=0.07,
                    text=f"• {bullet}",
                )
            )
        if recon.punchline:
            boxes.append(
                TextBoxNorm(
                    reading_order=900,
                    left=0.05,
                    top=0.88,
                    width=0.9,
                    height=0.08,
                    text=recon.punchline,
                    font_emphasis="small_caption",
                )
            )

    sub_rgb = RGBColor(*theme_colors["subtext"])
    for tb in boxes:
        if not (tb.text or "").strip():
            continue
        left_i, top_i, width_i, height_i = _norm_rect_to_inches(tb)
        try:
            shape = slide.shapes.add_textbox(left_i, top_i, width_i, height_i)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.text = tb.text.strip()
            p = tf.paragraphs[0]
            if tb.font_emphasis == "bold":
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = text_color
            elif tb.font_emphasis == "small_caption":
                p.font.size = Pt(12)
                p.font.italic = True
                p.font.color.rgb = sub_rgb
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = text_color
            _apply_aptos_narrow(shape, font_color=p.font.color.rgb)
        except Exception as e:
            logger.warning(f"Textbox add failed: {e}")

    return prs


def image_to_presentation(
    image_source: str,
    is_url: bool = True,
    webhook_url: str = None,
    api_key: str = "",
    layout_theme: str = DEFAULT_LAYOUT_THEME,
) -> dict:
    """
    Image → rich text/layout analysis → editable PPTX with textboxes, AI icons, optional source image panel.
    Without API keys, falls back to a single slide with the image fitted.
    """
    stats["requests_received"] += 1
    stats["last_request_time"] = datetime.now().isoformat()

    execution_id = str(uuid.uuid4())
    run_dir = os.path.join(OUTPUT_DIR, execution_id)
    os.makedirs(run_dir, exist_ok=True)

    try:
        if is_url:
            if image_source.startswith(("http://", "https://")):
                response = requests.get(image_source, verify=False)
                response.raise_for_status()
                img = Image.open(BytesIO(response.content))
            else:
                if ";base64," in image_source:
                    _, b64_data = image_source.split(";base64,", 1)
                    img = Image.open(BytesIO(base64.b64decode(b64_data)))
                else:
                    img = Image.open(image_source)
        else:
            img = Image.open(BytesIO(base64.b64decode(image_source)))

        img_ext = (img.format or "PNG").lower()
        if img_ext == "jpeg":
            img_ext = "jpg"
        img_path = os.path.join(run_dir, f"source_image.{img_ext}")
        if img.mode == "RGBA" and img_ext in ("jpg", "jpeg"):
            img = img.convert("RGB")
        img.save(img_path)

        # Vision APIs: prefer PNG for Anthropic
        png_path = os.path.join(run_dir, "source_image_for_vision.png")
        img.convert("RGBA" if img.mode in ("RGBA", "P") else "RGB").save(png_path, "PNG")
        pil_for_gemini = Image.open(png_path)

        has_genai = False
        has_anthropic = False
        client = None
        use_anthropic = False
        if api_key.startswith("sk-ant") or (
            not api_key
            and os.environ.get("ANTHROPIC_API_KEY")
        ):
            use_anthropic = True
        if use_anthropic:
            try:
                k = api_key if api_key else os.environ.get("ANTHROPIC_API_KEY")
                proxy_url = os.environ.get("GCP_PROXY_FOR_CLAUD")
                if proxy_url:
                    client = anthropic.Anthropic(api_key=k, base_url=proxy_url, max_retries=0)
                else:
                    client = anthropic.Anthropic(api_key=k, max_retries=0)
                has_anthropic = True
            except Exception:
                pass
        else:
            try:
                if api_key:
                    client = genai.Client(api_key=api_key)
                elif os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY"):
                    if not os.environ.get("GCP_PROXY_FOR_CLAUD"):
                        client = genai.Client()
                if client:
                    has_genai = True
            except Exception:
                pass

        if not (has_genai or has_anthropic):
            _send_progress(webhook_url, "No AI client; embedding image only...")
            prs = Presentation()
            prs.slide_width = SLIDE_WIDTH
            prs.slide_height = SLIDE_HEIGHT
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            img_width, img_height = img.size
            page_aspect = img_width / img_height
            slide_aspect = (SLIDE_WIDTH - 2 * MARGIN) / (SLIDE_HEIGHT - 2 * MARGIN)
            if page_aspect > slide_aspect:
                width = SLIDE_WIDTH - 2 * MARGIN
                height = width / page_aspect
            else:
                height = SLIDE_HEIGHT - 2 * MARGIN
                width = height * page_aspect
            left = (SLIDE_WIDTH - width) / 2 + MARGIN
            top = (SLIDE_HEIGHT - height) / 2 + MARGIN
            slide.shapes.add_picture(img_path, left, top, width, height)
            output_filename = "image_presentation.pptx"
            output_path = os.path.join(run_dir, output_filename)
            prs.save(output_path)
            file_url = _get_file_url(execution_id, output_filename)
            stats["successful_creations"] += 1
            _add_to_history(execution_id, output_filename, file_url, "image_to_pptx")
            response_payload = {
                "success": True,
                "message": "Image presentation generated (no API key: image embedded only). Add GEMINI_API_KEY or ANTHROPIC_API_KEY for full text + layout reconstruction.",
                "file_url": file_url,
                "download_path": f"/downloads/{execution_id}/{output_filename}",
                "execution_id": execution_id,
                "filename": output_filename,
            }
            _trigger_webhook(webhook_url, response_payload)
            return response_payload

        # --- Phase 1: meticulous image → text + layout description ---
        _send_progress(webhook_url, "Phase 1: extracting text and describing layout from image...")
        phase1_prompt = """You are an expert document and slide analyst. Perform a meticulous image-to-text and layout analysis.

CRITICAL:
1) Transcribe ALL visible text into full_text_reading_order and extracted_text_blocks (accurate wording; note region_hint for each block: e.g. top_title, left_column, diagram_label, legend, footer).
2) layout_description: describe spatial layout — columns, headers, spacing, alignment, visual hierarchy.
3) diagram_structure: describe diagrams faithfully — boxes, arrows, flows, charts, connectors, groupings (even if approximate).
4) color_and_style_notes: bullets on emphasis, colors, photos vs drawings.
5) visual_motifs_for_icons: 5-12 short English keywords for symbols/metaphors visible (used later for AI-generated icons matching the image).

Be exhaustive; downstream steps rebuild editable slides from this analysis."""

        if has_anthropic:
            with open(png_path, "rb") as f:
                b64_img = base64.b64encode(f.read()).decode("utf-8")
            raw_p1 = _call_anthropic_with_retry(b64_img, phase1_prompt, ImageLayoutAnalysisPhase1)
        else:
            resp_p1 = _call_genai_with_retry(client, pil_for_gemini, phase1_prompt, ImageLayoutAnalysisPhase1)
            raw_p1 = resp_p1.text.strip()

        if raw_p1.startswith("```json"):
            raw_p1 = raw_p1[7:]
        elif raw_p1.startswith("```"):
            raw_p1 = raw_p1[3:]
        if raw_p1.rstrip().endswith("```"):
            raw_p1 = raw_p1.rstrip()[:-3]
        phase1 = ImageLayoutAnalysisPhase1.model_validate(json.loads(raw_p1.strip()))

        # --- Phase 2: structured editable slide + icon seeds + text box geometry ---
        _send_progress(webhook_url, "Phase 2: reconstructing editable slide, text boxes, and icon seeds...")
        phase2_prompt = f"""You convert a completed Phase-1 image analysis into a structured editable PowerPoint specification.

Rules:
- Preserve ALL meaningful text from the analysis. Put verbatim copy into text_boxes with normalized coordinates (0-1 for full slide width/height, origin top-left). Approximate where each block sat in the image.
- Also set title, narrative, punchline, bullet_points for a coherent speaker-friendly summary (content must remain consistent with extracted_text_blocks).
- icon_keyword: best single seed for the main metaphor; extra_icon_keywords: more seeds from visual_motifs_for_icons (for DiceBear icons).
- place_original_image_as_reference: true if charts, complex diagrams, or precise geometry must remain visible; then set reference_image_box OR leave null for default right-hand panel (52%-98% horizontal).
- layout_type: diagram | two_column | title_and_content — semantic only (we use one slide with positioned boxes).

PHASE 1 ANALYSIS JSON:
{json.dumps(phase1.model_dump(), indent=2)}
"""
        recon_data = _llm_json_structured(client, use_anthropic, phase2_prompt, ImageToPptReconstruction)
        recon = ImageToPptReconstruction.model_validate(recon_data)

        _send_progress(webhook_url, "Building PowerPoint with text boxes and icons...")
        prs = _build_presentation_from_image_reconstruction(run_dir, img_path, recon, layout_theme)

        output_filename = "image_presentation.pptx"
        output_path = os.path.join(run_dir, output_filename)
        prs.save(output_path)

        file_url = _get_file_url(execution_id, output_filename)

        stats["successful_creations"] += 1
        _add_to_history(execution_id, output_filename, file_url, "image_to_pptx")

        response_payload = {
            "success": True,
            "message": "Image analyzed (text + layout), reconstructed with editable text boxes, AI icons, and optional source diagram panel.",
            "file_url": file_url,
            "download_path": f"/downloads/{execution_id}/{output_filename}",
            "execution_id": execution_id,
            "filename": output_filename,
        }
        _trigger_webhook(webhook_url, response_payload)
        return response_payload

    except Exception as e:
        stats["failed_creations"] += 1
        error_payload = {
            "success": False,
            "message": f"Error converting image to presentation: {str(e)}",
        }
        _trigger_webhook(webhook_url, error_payload)
        return error_payload
def generate_artifacts_from_prompt(
    prompt: str,
    target_format: str = "pptx",
    presentation_style: str = "Detailed",
    layout_theme: str = DEFAULT_LAYOUT_THEME,
    num_slides: int = 5,
    webhook_url: str = None,
    api_key: str = ""
) -> dict:
    """
    Dynamically generates a full presentation or document strictly from a text prompt.
    Takes into account requested themes and styles.
    """
    stats["requests_received"] += 1
    stats["last_request_time"] = datetime.now().isoformat()
    
    execution_id = str(uuid.uuid4())
    run_dir = os.path.join(OUTPUT_DIR, execution_id)
    os.makedirs(run_dir, exist_ok=True)
    
    try:
        has_genai = False
        has_anthropic = False
        client = None
        
        use_anthropic = False
        if api_key.startswith("sk-ant") or (not api_key and os.environ.get("ANTHROPIC_API_KEY")):
            use_anthropic = True
            
        if use_anthropic:
            try:
                k = api_key if api_key else os.environ.get("ANTHROPIC_API_KEY")
                proxy_url = os.environ.get("GCP_PROXY_FOR_CLAUD")
                if proxy_url:
                    client = anthropic.Anthropic(api_key=k, base_url=proxy_url, max_retries=0)
                else:
                    client = anthropic.Anthropic(api_key=k, max_retries=0)
                has_anthropic = True
            except:
                pass
        else:
            try:
                if api_key:
                    client = genai.Client(api_key=api_key)
                elif os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY"):
                    # DO NOT use Gemini if GCP proxy is configured (Anthropic only for proxy environment)
                    if not os.environ.get("GCP_PROXY_FOR_CLAUD"):
                        client = genai.Client()
                if client:
                    has_genai = True
            except:
                pass
                
        if not (has_genai or has_anthropic):
            raise Exception("No valid GenAI or Anthropic API key configured.")
            
        if target_format.lower() == "pptx":
            # Phase A: advance layout plan (every slide's layout and role planned before content)
            _send_progress(webhook_url, "Planning deck layout (story arc and per-slide layout)...")
            plan_prompt = f"""You are an expert presentation strategist. Plan a {num_slides}-slide deck on this topic:
{prompt}

Style: {presentation_style}. Theme: {layout_theme}.

Requirements:
- Produce EXACTLY {num_slides} entries in \"slides\", indices 0..{num_slides - 1} in order.
- Choose layout_type per slide: title_slide (opening), section_divider (chapter breaks), index_slide (agenda/TOC if needed), title_and_content, two_column, or diagram as appropriate.
- slide_archetype must match intent: title, agenda, divider, standard, table, deep_dive, roadmap, options.
- For each slide, state purpose_one_line and visual_role (how the AI icon will reinforce the message).
- Content slides are rendered as an infographic: one row per bullet with an AI icon — plan concise bullets (no long paragraphs).

Deck flow: Context/Vision -> Execution -> Options -> Architecture -> Roadmap -> Risks -> Recommendation where applicable.
Use neutral wording (no JPL/JEMP unless user asked)."""
            plan_dict = _llm_json_structured(client, use_anthropic, plan_prompt, PresentationLayoutPlan)
            if len(plan_dict.get("slides") or []) != num_slides:
                raise Exception(
                    f"Layout plan must contain exactly {num_slides} slides; got {len(plan_dict.get('slides') or [])}."
                )

            # Phase B: full slide JSON following the approved plan
            _send_progress(webhook_url, "Generating slide content from approved layout plan...")
            system_prompt = f"""You are an expert presentation designer and strategic consultant.
Create a {num_slides}-slide presentation on the following topic: {prompt}

Presentation Style Constraint: {presentation_style}
Theme Concept: {layout_theme}

APPROVED LAYOUT PLAN (you MUST follow slide order, layout_type, and archetype per index):
{json.dumps(plan_dict, indent=2)}

STRICT VALIDATION RULES (Generic Presentation Kit):
DO NOT brand with 'JPL', 'JEMP', or corporate tags unless the user asked. Use neutral terms.
Visuals: default themes render every content slide with AI-generated icons — one icon per bullet row plus a larger hero icon (DiceBear seeds from icon_keyword and bullet_icon_seeds). A two-column split layout is used only if the theme name includes phrases like "split layout" or "two-panel".

CONTENT SLIDES (all slides that are NOT title_slide, NOT section_divider, NOT index_slide, and NOT archetype title/agenda/divider) MUST EACH HAVE:
1. title — short headline only (about 10 words max, ~72 characters) so it stays on ONE line when rendered.
2. narrative — non-empty; at most TWO short lines (~220 characters total); must feel more prominent than bullets.
3. bullet_points — 3-5 bullets; EACH bullet max ~120 characters — terse, one idea per row (infographic layout).
4. bullet_icon_seeds — REQUIRED: same number of entries as bullet_points; each a short English seed for a matching AI icon beside that bullet.
5. icon_keyword — non-empty; also used for the large hero icon (top-right).
6. punchline — non-empty (one takeaway at the bottom).

NON-CONTENT SLIDES (title_slide, section_divider, index_slide, or archetypes title/agenda/divider): do not require the full five-part anatomy; keep them clean and readable.

Each content slide MUST also respect the planned layout_type from the layout plan.

DECK STRUCTURE:
- Storyline Flow: Context/Vision -> Execution Model -> Options -> Architecture -> Roadmap -> Risks -> Recommendation.
- For comparisons/options, use table_data where appropriate.

Output JSON matching the PresentationData schema (top-level key \"slides\" only)."""
            data = _llm_json_structured(client, use_anthropic, system_prompt, PresentationData)
            slides_data = data.get("slides") or []
            if len(slides_data) != num_slides:
                raise Exception(f"Expected {num_slides} slides; model returned {len(slides_data)}.")

            # Structural validation + LLM repair loop
            for attempt in range(3):
                errs = _validate_strict_content_slides(slides_data)
                if not errs:
                    break
                _send_progress(webhook_url, f"Repairing slide data (validation pass {attempt + 1}/3)...")
                slides_data = _repair_presentation_slides(
                    client, use_anthropic, slides_data, errs, plan_dict, None
                )
            final_errs = _validate_strict_content_slides(slides_data)
            if final_errs:
                raise Exception("Slide validation failed after repair: " + "; ".join(final_errs[:12]))

            # Two mandatory AI visual-layout QA rounds (overlap / crowding / balance vs. fixed regions)
            _send_progress(webhook_url, "Visual QA review round 1 (layout and overlap risk)...")
            vqa1 = _run_visual_layout_review(client, use_anthropic, slides_data, plan_dict, 1, None)
            _send_progress(webhook_url, "Visual QA review round 2 (layout and overlap risk)...")
            vqa2 = _run_visual_layout_review(client, use_anthropic, slides_data, plan_dict, 2, vqa1)
            qa_notes = []
            for r, tag in ((vqa1, "R1"), (vqa2, "R2")):
                for it in r.issues:
                    qa_notes.append(f"[{tag}] Slide {it.slide_index} ({it.severity}): {it.issue}")
            qa_notes.extend(vqa2.suggested_fixes or [])
            if qa_notes:
                _send_progress(webhook_url, "Applying visual QA feedback to slide copy...")
                slides_data = _repair_presentation_slides(
                    client, use_anthropic, slides_data, [], plan_dict, qa_notes
                )
                fe = _validate_strict_content_slides(slides_data)
                if fe:
                    slides_data = _repair_presentation_slides(
                        client, use_anthropic, slides_data, fe, plan_dict, None
                    )

            output_filename = "generated_presentation.pptx"
            output_path = os.path.join(run_dir, output_filename)
            
            prs, theme_colors = _create_themed_presentation(layout_theme)
            bg_color = RGBColor(*theme_colors["bg"])
            text_color = RGBColor(*theme_colors["text"])
            theme_low = layout_theme.lower()
            title_accent_color = _title_color_from_theme(theme_colors, text_color)
            large_title_pt = 32 if _is_dark_studio_theme(layout_theme) else None

            _send_progress(webhook_url, "Building presentation file...")

            for i, s_data in enumerate(slides_data):
                l_type = s_data.get("layout_type", "title_and_content")
                is_strict = _is_strict_content_slide(s_data)

                if l_type == "title_slide":
                    slide_layout = prs.slide_layouts[0]
                elif l_type == "section_divider":
                    slide_layout = prs.slide_layouts[2]
                elif l_type == "index_slide":
                    slide_layout = prs.slide_layouts[1]
                elif is_strict:
                    slide_layout = prs.slide_layouts[6]
                elif l_type == "two_column":
                    slide_layout = prs.slide_layouts[3]
                else:
                    slide_layout = prs.slide_layouts[1]

                slide = prs.slides.add_slide(slide_layout)
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = bg_color
                _apply_theme_ribbons(slide, prs, theme_colors)

                if l_type == "title_slide":
                    if slide.shapes.title:
                        _style_slide_title_shape(
                            slide.shapes.title,
                            s_data.get("title", f"Slide {i + 1}"),
                            title_accent_color,
                            truncate=True,
                            font_pt=large_title_pt,
                        )
                    if len(slide.placeholders) > 1:
                        subtitle_shape = slide.placeholders[1]
                        subtitle_shape.text = (s_data.get("narrative") or s_data.get("punchline") or "")[:650]
                        for p in subtitle_shape.text_frame.paragraphs:
                            p.font.size = Pt(DECK_NARRATIVE_PT)
                            p.alignment = PP_ALIGN.LEFT
                        _apply_aptos_narrow(subtitle_shape, font_color=RGBColor(*theme_colors["subtext"]))
                    continue

                if l_type == "section_divider":
                    if slide.shapes.title:
                        _style_slide_title_shape(
                            slide.shapes.title,
                            s_data.get("title", "Section"),
                            title_accent_color,
                            truncate=True,
                            font_pt=large_title_pt,
                        )
                    if len(slide.placeholders) > 1:
                        subtitle_shape = slide.placeholders[1]
                        subtitle_shape.text = (s_data.get("narrative") or s_data.get("punchline") or "")[:650]
                        for p in subtitle_shape.text_frame.paragraphs:
                            p.font.size = Pt(DECK_NARRATIVE_PT)
                            p.alignment = PP_ALIGN.LEFT
                        _apply_aptos_narrow(subtitle_shape, font_color=RGBColor(*theme_colors["subtext"]))
                    continue

                if l_type == "index_slide":
                    if slide.shapes.title:
                        _style_slide_title_shape(
                            slide.shapes.title,
                            s_data.get("title", "Agenda"),
                            title_accent_color,
                            truncate=True,
                            font_pt=large_title_pt,
                        )
                    if len(slide.placeholders) > 1:
                        body_shape = slide.placeholders[1]
                        body_shape.left = Inches(0.5)
                        body_shape.top = Inches(1.35)
                        body_shape.width = SLIDE_WIDTH - Inches(1.0)
                        body_shape.height = Inches(5.65)
                        tf = body_shape.text_frame
                        tf.word_wrap = True
                        tf.text = ""
                        for bullet in s_data.get("bullet_points", []) or []:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                            p.font.size = Pt(DECK_BULLET_PT)
                        _apply_aptos_narrow(body_shape, font_color=text_color)
                    continue

                if is_strict:
                    _add_strict_content_slide_infographic(
                        slide,
                        s_data,
                        run_dir,
                        i,
                        theme_colors,
                        text_color,
                        theme_low,
                        layout_theme,
                    )
                    continue

                # Rare fallback: non-strict content slides (placeholder layout + hierarchy)
                narrative_width = SLIDE_WIDTH - Inches(1.0) - Inches(1.15)
                if slide.shapes.title:
                    ts = slide.shapes.title
                    ts.left = Inches(0.5)
                    ts.top = Inches(0.25)
                    ts.width = narrative_width
                    ts.height = Inches(0.8)
                    _style_slide_title_shape(
                        ts,
                        s_data.get("title", f"Slide {i + 1}"),
                        title_accent_color,
                        truncate=True,
                        font_pt=large_title_pt,
                    )

                tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), narrative_width, Inches(0.5))
                ntf = tx_box.text_frame
                ntf.word_wrap = True
                np = ntf.paragraphs[0]
                np.text = s_data.get("narrative", "")
                np.font.size = Pt(DECK_NARRATIVE_PT)
                np.font.color.rgb = text_color
                np.alignment = PP_ALIGN.LEFT
                _apply_aptos_narrow(tx_box, font_color=text_color)

                punch = slide.shapes.add_textbox(Inches(0.5), SLIDE_HEIGHT - Inches(0.8), SLIDE_WIDTH - Inches(1.0), Inches(0.4))
                pp = punch.text_frame.paragraphs[0]
                pp.text = s_data.get("punchline", "")
                pp.font.italic = True
                pp.font.size = Pt(DECK_PUNCHLINE_PT)
                pp.font.color.rgb = RGBColor(*theme_colors["subtext"])
                pp.alignment = PP_ALIGN.LEFT
                _apply_aptos_narrow(punch, font_color=pp.font.color.rgb)

                if len(slide.placeholders) > 1:
                    body_shape = slide.placeholders[1]
                    body_shape.left = Inches(0.5)
                    body_shape.top = Inches(1.6)
                    body_shape.width = narrative_width
                    body_shape.height = Inches(4.8)
                    tf = body_shape.text_frame
                    tf.word_wrap = True
                    tf.text = ""
                    for bullet in s_data.get("bullet_points", []):
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(DECK_BULLET_PT)
                    _apply_aptos_narrow(body_shape, font_color=text_color)

                ik = (s_data.get("icon_keyword") or "presentation").strip()
                ipath = os.path.join(run_dir, f"icon_fallback_{i}.png")
                if _download_dicebear_icon(ik, ipath):
                    slide.shapes.add_picture(ipath, Inches(11.0), Inches(0.35), Inches(1.2), Inches(1.2))

                if l_type == "two_column" and len(slide.placeholders) > 2:
                    body_shape.width = (SLIDE_WIDTH / 2) - Inches(0.75)
                    right_body_shape = slide.placeholders[2]
                    right_body_shape.left = (SLIDE_WIDTH / 2) + Inches(0.25)
                    right_body_shape.top = Inches(1.6)
                    right_body_shape.width = (SLIDE_WIDTH / 2) - Inches(0.75)
                    right_body_shape.height = Inches(4.8)
                    tf_right = right_body_shape.text_frame
                    tf_right.word_wrap = True
                    tf_right.text = "Additional Context / Visuals"
                    _apply_aptos_narrow(right_body_shape, font_color=text_color)

            prs.save(output_path)
            
        else: # DOCX
            _send_progress(webhook_url, "Generating document content with AI...")
            system_prompt = f"""You are an expert document author and strategic consultant.
Create a detailed, well-structured document on the following topic: {prompt}

Document Style Constraint: {presentation_style}
Theme/Tone: {layout_theme}

STRICT DOCUMENT ORGANIZATION (The Formatting Kit):
1. Introduction: Must explain what the document is and how sections are organized. Outline the roadmap at a section level.
2. Structure: Follow a logical order (e.g., Vision -> Execution/Tracks -> Architecture/Options -> Details).
3. Clarity: Separate Projects (what you build) from Tracks (how you execute). Include comparative analysis where options are discussed.
4. Formatting: The document must be well organized into headings, content paragraphs, and bullet points where useful for scanning.
5. Depth: Do not drop content, only add. Avoid padding just to hit a page count.
"""
            if use_anthropic:
                raw_text = _call_anthropic_text_with_retry(client, system_prompt, DocumentData)
            else:
                response = _call_genai_text_with_retry(client, system_prompt, DocumentData)
                raw_text = response.text
                
            try:
                raw_text = raw_text.strip()
                if raw_text.startswith("```json"): raw_text = raw_text[7:]
                elif raw_text.startswith("```"): raw_text = raw_text[3:]
                if raw_text.endswith("```"): raw_text = raw_text[:-3]
                data = json.loads(raw_text.strip())
            except Exception as e:
                raise Exception(f"Failed to parse LLM JSON: {e}")
                
            output_filename = "generated_document.docx"
            output_path = os.path.join(run_dir, output_filename)
            
            docx_doc = DocxDocument()
            docx_doc.add_heading(data.get("title", "Generated Document"), 0)
            
            for section in data.get("sections", []):
                docx_doc.add_heading(section.get("heading", "Section"), level=1)
                for paragraph in section.get("content", "").split("\n\n"):
                    if paragraph.strip():
                        docx_doc.add_paragraph(paragraph.strip())
                for bullet in section.get("bullet_points", []):
                    docx_doc.add_paragraph(bullet, style='List Bullet')
                    
            docx_doc.save(output_path)
            
            # Apply formatting guidelines
            formatted_output_filename = "final_formatted_document.docx"
            formatted_output_path = os.path.join(run_dir, formatted_output_filename)
            apply_guidelines(output_path, formatted_output_path)
            output_filename = formatted_output_filename
            
        file_url = _get_file_url(execution_id, output_filename)
        stats["successful_creations"] += 1
        _add_to_history(execution_id, output_filename, file_url, "generate_from_prompt")
        
        response_payload = {
            "success": True,
            "message": f"Successfully generated {target_format.upper()} from prompt.",
            "file_url": file_url,
            "download_path": f"/downloads/{execution_id}/{output_filename}",
            "execution_id": execution_id,
            "filename": output_filename
        }
        _trigger_webhook(webhook_url, response_payload)
        return response_payload

    except Exception as e:
        stats["failed_creations"] += 1
        error_payload = {
            "success": False,
            "message": f"Error generating from prompt: {str(e)}"
        }
        _trigger_webhook(webhook_url, error_payload)
        return error_payload

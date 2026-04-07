import os
import shutil
import tempfile
import base64

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import HTMLResponse, FileResponse, Response
from pydantic import BaseModel
import uvicorn

from core import generate_presentation, image_to_presentation, format_document, process_pdf_to_artifacts, stats, generation_history, OUTPUT_DIR
from mcp_server import mcp
from werkzeug.middleware.proxy_fix import ProxyFix
from starlette.middleware.wsgi import WSGIMiddleware

app = FastAPI(title="PPTX Generator API", description="API and UI for generating PowerPoint presentations")

if os.environ.get("GCP_PROXY_FOR_CLAUD"):
    from fastapi.middleware.trustedhost import TrustedHostMiddleware
    app.add_middleware(TrustedHostMiddleware, allowed_hosts=["*"])
    from starlette.middleware.base import BaseHTTPMiddleware
    
    class ProxyHeadersMiddleware(BaseHTTPMiddleware):
        async def dispatch(self, request, call_next):
            # Trust downstream proxies (similar to ProxyFix)
            request.scope["scheme"] = request.headers.get("x-forwarded-proto", request.scope.get("scheme", "http"))
            return await call_next(request)
            
    app.add_middleware(ProxyHeadersMiddleware)

MAX_UPLOAD_SIZE_BYTES = int(os.environ.get("MAX_UPLOAD_SIZE_BYTES", str(25 * 1024 * 1024)))

# Mount the MCP SSE application
mcp_starlette = mcp.sse_app()
app.mount("/mcp", mcp_starlette)

# Pydantic models for API
class GenerateRequest(BaseModel):
    python_code: str
    webhook_url: str = None

class ImageRequest(BaseModel):
    image_source: str
    is_url: bool = True
    webhook_url: str = None

class DocxRequest(BaseModel):
    doc_source: str
    is_url: bool = True
    webhook_url: str = None

class PdfRequest(BaseModel):
    pdf_source: str
    is_url: bool = True
    instructions: str = ""
    layout_theme: str = ""
    visual_iconography: str = ""
    slide_content_rules: str = ""
    target_format: str = "pptx"
    webhook_url: str = None
    api_key: str = ""


def _persist_upload_to_tempfile(upload: UploadFile, suffix: str) -> str:
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    try:
        shutil.copyfileobj(upload.file, temp_file)
    finally:
        temp_file.close()
    return temp_file.name


def _persist_upload_to_tempfile(upload: UploadFile, suffix: str) -> str:
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    try:
        shutil.copyfileobj(upload.file, temp_file)
    finally:
        temp_file.close()
    return temp_file.name

@app.get("/", response_class=HTMLResponse)
async def index():
    html_content = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>PPTX Generator Dashboard</title>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 40px; background-color: #f5f5f5; }}
            .container {{ max-width: 800px; margin: 0 auto; background: white; padding: 20px; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); }}
            h1 {{ color: #333; }}
            .stats {{ display: flex; gap: 20px; margin-bottom: 30px; }}
            .stat-box {{ flex: 1; background: #e9ecef; padding: 15px; border-radius: 6px; text-align: center; }}
            .stat-box h3 {{ margin: 0 0 10px 0; font-size: 14px; color: #666; }}
            .stat-box p {{ margin: 0; font-size: 24px; font-weight: bold; color: #333; }}
            textarea {{ width: 100%; height: 200px; margin-bottom: 15px; padding: 10px; font-family: monospace; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box; }}
            button {{ background: #007bff; color: white; border: none; padding: 10px 20px; border-radius: 4px; cursor: pointer; font-size: 16px; }}
            button:hover {{ background: #0056b3; }}
            .result {{ margin-top: 20px; padding: 15px; border-radius: 4px; display: none; }}
            .success {{ background: #d4edda; color: #155724; border: 1px solid #c3e6cb; }}
            .error {{ background: #f8d7da; color: #721c24; border: 1px solid #f5c6cb; }}
            .history-list {{ list-style-type: none; padding: 0; }}
            .history-list li {{ background: #f8f9fa; margin-bottom: 10px; padding: 15px; border-radius: 6px; border-left: 4px solid #007bff; display: flex; justify-content: space-between; align-items: center; }}
            .history-list li a {{ color: #007bff; text-decoration: none; font-weight: bold; }}
            .history-list li a:hover {{ text-decoration: underline; }}
            .history-meta {{ font-size: 12px; color: #6c757d; }}
        </style>
    </head>
    <body>
        <div class="container">
            <h1>PPTX Generator Dashboard</h1>
            
            <div class="stats">
                <div class="stat-box">
                    <h3>Requests Received</h3>
                    <p>{stats['requests_received']}</p>
                </div>
                <div class="stat-box">
                    <h3>Successful Creations</h3>
                    <p>
                        {
                            f'<a href="{stats["last_success_file_url"]}" title="Download {stats["last_success_filename"] or "latest file"}" style="color: inherit; text-decoration: underline;">{stats["successful_creations"]}</a>'
                            if stats.get("last_success_file_url")
                            else str(stats["successful_creations"])
                        }
                    </p>
                </div>
                <div class="stat-box">
                    <h3>Failed Creations</h3>
                    <p>{stats['failed_creations']}</p>
                </div>
            </div>

            <div style="display: flex; gap: 10px; margin-bottom: 20px;">
                <button id="tabCode" style="flex: 1; background: #007bff;" onclick="switchTab('code')">Generate from Code</button>
                <button id="tabImage" style="flex: 1; background: #6c757d;" onclick="switchTab('image')">Image to PPTX</button>
                <button id="tabDocx" style="flex: 1; background: #6c757d;" onclick="switchTab('docx')">Format DOCX Template</button>
                <button id="tabPdf" style="flex: 1; background: #6c757d;" onclick="switchTab('pdf')">Process PDF</button>
            </div>

            <div id="sectionCode">
                <h2>Trigger Presentation Creation</h2>
                <form id="generateForm">
                    <textarea id="pythonCode" placeholder="Enter python-pptx code here... Example:
from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Hello World'
prs.save('output.pptx')"></textarea>
                    <button type="submit" id="submitBtn">Generate PPTX</button>
                </form>
            </div>

            <div id="sectionImage" style="display: none;">
                <h2>Convert Image to PPTX</h2>
                <form id="imageForm">
                    <div style="margin-bottom: 15px;">
                        <label><strong>Option 1: Upload Image</strong></label>
                        <input type="file" id="imageFileInput" accept="image/*" style="display: block; margin-top: 5px;">
                    </div>
                    <div style="text-align: center; margin-bottom: 15px;">OR</div>
                    <label><strong>Option 2: Provide URL</strong></label>
                    <input type="text" id="imageSource" placeholder="Enter Image URL..." style="width: 100%; padding: 10px; margin-top: 5px; margin-bottom: 15px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;">
                    <button type="submit" id="submitImageBtn">Convert to PPTX</button>
                </form>
            </div>

            <div id="sectionDocx" style="display: none;">
                <h2>Format Document (Apply Template)</h2>
                <form id="docxForm">
                    <div style="margin-bottom: 15px;">
                        <label><strong>Option 1: Upload DOCX</strong></label>
                        <input type="file" id="docxFileInput" accept=".docx" style="display: block; margin-top: 5px;">
                    </div>
                    <div style="text-align: center; margin-bottom: 15px;">OR</div>
                    <label><strong>Option 2: Provide URL</strong></label>
                    <input type="text" id="docxSource" placeholder="Enter DOCX File URL..." style="width: 100%; padding: 10px; margin-top: 5px; margin-bottom: 15px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;">
                    <button type="submit" id="submitDocxBtn">Format Document</button>
                </form>
                <div style="margin: 12px 0; text-align: center; color: #666;">or</div>
                <form id="docxUploadForm">
                    <input type="file" id="docxFile" accept=".docx" style="width: 100%; margin-bottom: 15px;" required>
                    <button type="submit" id="submitDocxUploadBtn">Upload DOCX & Format</button>
                </form>
            </div>

            <div id="sectionPdf" style="display: none;">
                <h2>Process PDF to Presentation/Document</h2>
                <form id="pdfForm">
                    <div style="margin-bottom: 15px;">
                        <label><strong>Option 1: Upload PDF</strong></label>
                        <input type="file" id="pdfFileInput" accept=".pdf" style="display: block; margin-top: 5px;">
                    </div>
                    <div style="text-align: center; margin-bottom: 15px;">OR</div>
                    <label><strong>Option 2: Provide URL</strong></label>
                    <input type="text" id="pdfSource" placeholder="Enter PDF File URL..." style="width: 100%; padding: 10px; margin-top: 5px; margin-bottom: 15px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;">
                    
                    <textarea id="pdfInstructions" placeholder="Abstract Instructions (e.g., 'Extract financial tables only')..." style="height: 60px;"></textarea>
                    <input type="text" id="pdfLayoutTheme" placeholder="Layout Theme (e.g., 'Modern Corporate')" style="width: 100%; padding: 10px; margin-bottom: 10px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;">
                    <input type="text" id="pdfIconography" placeholder="Visual Iconography (e.g., 'Flat design, tech icons')" style="width: 100%; padding: 10px; margin-bottom: 10px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;">
                    <textarea id="pdfContentRules" placeholder="Slide Content Rules (e.g., 'Max 5 bullets per slide')..." style="height: 60px;"></textarea>
                    
                    <input type="text" id="pdfApiKey" placeholder="API Key (optional, defaults to server env. Gemini or Anthropic allowed)" style="width: 100%; padding: 10px; margin-bottom: 15px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;" />

                    <select id="pdfTargetFormat" style="width: 100%; padding: 10px; margin-bottom: 15px; border: 1px solid #ccc; border-radius: 4px;">
                        <option value="pptx">Generate PPTX</option>
                        <option value="docx">Generate DOCX</option>
                    </select>

                    <button type="submit" id="submitPdfBtn">Process PDF</button>
                </form>
                <div style="margin: 12px 0; text-align: center; color: #666;">or</div>
                <form id="pdfUploadForm">
                    <input type="file" id="pdfFile" accept=".pdf" style="width: 100%; margin-bottom: 15px;" required>
                    <textarea id="pdfUploadInstructions" placeholder="Abstract Instructions..." style="height: 60px;"></textarea>
                    <input type="text" id="pdfUploadLayoutTheme" placeholder="Layout Theme" style="width: 100%; padding: 10px; margin-bottom: 10px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;">
                    <input type="text" id="pdfUploadIconography" placeholder="Visual Iconography" style="width: 100%; padding: 10px; margin-bottom: 10px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;">
                    <textarea id="pdfUploadContentRules" placeholder="Slide Content Rules..." style="height: 60px;"></textarea>
                    <input type="text" id="pdfUploadApiKey" placeholder="API Key (optional, defaults to server env. Gemini or Anthropic allowed)" style="width: 100%; padding: 10px; margin-bottom: 15px; border: 1px solid #ccc; border-radius: 4px; box-sizing: border-box;" />
                    <select id="pdfUploadTargetFormat" style="width: 100%; padding: 10px; margin-bottom: 15px; border: 1px solid #ccc; border-radius: 4px;">
                        <option value="pptx">Generate PPTX</option>
                        <option value="docx">Generate DOCX</option>
                    </select>
                    <button type="submit" id="submitPdfUploadBtn">Upload PDF & Process</button>
                </form>
            </div>

            <div id="resultBox" class="result"></div>
            
            <hr style="margin: 40px 0; border: 0; border-top: 1px solid #eee;">
            
            <h2>Recent Generations</h2>
            <div id="historyContainer">
                <p style="color: #666; font-style: italic;">Loading history...</p>
            </div>
        </div>

        <script>
            function loadHistory() {{
                fetch('/api/history')
                    .then(response => response.json())
                    .then(data => {{
                        const container = document.getElementById('historyContainer');
                        if (!data || data.length === 0) {{
                            container.innerHTML = '<p style="color: #666; font-style: italic;">No recent generations found.</p>';
                            return;
                        }}
                        
                        let html = '<ul class="history-list">';
                        data.forEach(item => {{
                            const date = new Date(item.timestamp).toLocaleString();
                            html += `
                                <li>
                                    <div>
                                        <a href="${{item.file_url}}" target="_blank" download>${{item.filename}}</a>
                                        <div class="history-meta">ID: ${{item.execution_id}} | Type: ${{item.type}}</div>
                                    </div>
                                    <div style="font-size: 12px; color: #999;">${{date}}</div>
                                </li>
                            `;
                        }});
                        html += '</ul>';
                        container.innerHTML = html;
                    }})
                    .catch(err => {{
                        console.error('Error loading history:', err);
                        document.getElementById('historyContainer').innerHTML = '<p style="color: red;">Failed to load history.</p>';
                    }});
            }}

            document.addEventListener('DOMContentLoaded', loadHistory);

            function fileToBase64(file) {{
                return new Promise((resolve, reject) => {{
                    const reader = new FileReader();
                    reader.readAsDataURL(file);
                    reader.onload = () => resolve(reader.result);
                    reader.onerror = error => reject(error);
                }});
            }}

            function switchTab(tab) {{
                document.getElementById('sectionCode').style.display = tab === 'code' ? 'block' : 'none';
                document.getElementById('sectionImage').style.display = tab === 'image' ? 'block' : 'none';
                document.getElementById('sectionDocx').style.display = tab === 'docx' ? 'block' : 'none';
                document.getElementById('sectionPdf').style.display = tab === 'pdf' ? 'block' : 'none';
                document.getElementById('tabCode').style.background = tab === 'code' ? '#007bff' : '#6c757d';
                document.getElementById('tabImage').style.background = tab === 'image' ? '#007bff' : '#6c757d';
                document.getElementById('tabDocx').style.background = tab === 'docx' ? '#007bff' : '#6c757d';
                document.getElementById('tabPdf').style.background = tab === 'pdf' ? '#007bff' : '#6c757d';
                document.getElementById('resultBox').style.display = 'none';
            }}

            document.getElementById('generateForm').addEventListener('submit', async (e) => {{
                e.preventDefault();
                const btn = document.getElementById('submitBtn');
                const resultBox = document.getElementById('resultBox');
                const code = document.getElementById('pythonCode').value;
                
                if (!code) return;
                
                btn.disabled = true;
                btn.textContent = 'Generating...';
                resultBox.style.display = 'none';
                
                try {{
                    const response = await fetch('/api/generate', {{
                        method: 'POST',
                        headers: {{ 'Content-Type': 'application/json' }},
                        body: JSON.stringify({{ python_code: code }})
                    }});
                    
                    const data = await response.json();
                    
                    resultBox.style.display = 'block';
                    if (data.success) {{
                        resultBox.className = 'result success';
                        resultBox.innerHTML = `<strong>Success!</strong> Presentation generated. <br><a href="${{data.file_url}}" target="_blank">Download ${{data.filename || 'File'}}</a>`;
                        
                        setTimeout(() => window.location.reload(), 2000);
                    }} else {{
                        resultBox.className = 'result error';
                        resultBox.innerHTML = `<strong>Error!</strong><br><pre>${{data.message}}</pre>`;
                    }}
                }} catch (err) {{
                    resultBox.style.display = 'block';
                    resultBox.className = 'result error';
                    resultBox.textContent = 'Network error occurred.';
                }} finally {{
                    btn.disabled = false;
                    btn.textContent = 'Generate PPTX';
                }}
            }});

            document.getElementById('imageForm').addEventListener('submit', async (e) => {{
                e.preventDefault();
                const btn = document.getElementById('submitImageBtn');
                const resultBox = document.getElementById('resultBox');
                const fileInput = document.getElementById('imageFileInput');
                let source = document.getElementById('imageSource').value;
                let is_url = true;

                if (fileInput.files.length > 0) {{
                    source = await fileToBase64(fileInput.files[0]);
                    is_url = true; // The backend naturally handles data URIs when is_url is true
                }}
                
                if (!source) {{
                    alert("Please provide an image file or URL.");
                    return;
                }}
                
                btn.disabled = true;
                btn.textContent = 'Converting...';
                resultBox.style.display = 'none';
                
                try {{
                    const response = await fetch('/api/image-to-pptx', {{
                        method: 'POST',
                        headers: {{ 'Content-Type': 'application/json' }},
                        body: JSON.stringify({{ image_source: source, is_url: true }})
                    }});
                    
                    const data = await response.json();
                    
                    resultBox.style.display = 'block';
                    if (data.success) {{
                        resultBox.className = 'result success';
                        resultBox.innerHTML = `<strong>Success!</strong> Presentation generated. <br><a href="${{data.file_url}}" target="_blank">Download ${{data.filename || 'File'}}</a>`;
                        
                        setTimeout(() => window.location.reload(), 2000);
                    }} else {{
                        resultBox.className = 'result error';
                        resultBox.innerHTML = `<strong>Error!</strong><br><pre>${{data.message}}</pre>`;
                    }}
                }} catch (err) {{
                    resultBox.style.display = 'block';
                    resultBox.className = 'result error';
                    resultBox.textContent = 'Network error occurred.';
                }} finally {{
                    btn.disabled = false;
                    btn.textContent = 'Convert to PPTX';
                }}
            }});
            document.getElementById('docxForm').addEventListener('submit', async (e) => {{
                e.preventDefault();
                const btn = document.getElementById('submitDocxBtn');
                const resultBox = document.getElementById('resultBox');
                const fileInput = document.getElementById('docxFileInput');
                let source = document.getElementById('docxSource').value;
                let is_url = true;

                if (fileInput.files.length > 0) {{
                    source = await fileToBase64(fileInput.files[0]);
                }}
                
                if (!source) {{
                    alert("Please provide a DOCX file or URL.");
                    return;
                }}
                
                btn.disabled = true;
                btn.textContent = 'Formatting...';
                resultBox.style.display = 'none';
                
                try {{
                    const response = await fetch('/api/format-docx', {{
                        method: 'POST',
                        headers: {{ 'Content-Type': 'application/json' }},
                        body: JSON.stringify({{ doc_source: source, is_url: true }})
                    }});
                    
                    const data = await response.json();
                    
                    resultBox.style.display = 'block';
                    if (data.success) {{
                        resultBox.className = 'result success';
                        resultBox.innerHTML = `<strong>Success!</strong> Document formatted. <br><a href="${{data.file_url}}" target="_blank">Download ${{data.filename || 'File'}}</a>`;
                        
                        setTimeout(() => window.location.reload(), 2000);
                    }} else {{
                        resultBox.className = 'result error';
                        resultBox.innerHTML = `<strong>Error!</strong><br><pre>${{data.message}}</pre>`;
                    }}
                }} catch (err) {{
                    resultBox.style.display = 'block';
                    resultBox.className = 'result error';
                    resultBox.textContent = 'Network error occurred.';
                }} finally {{
                    btn.disabled = false;
                    btn.textContent = 'Format Document';
                }}
            }});
            document.getElementById('docxUploadForm').addEventListener('submit', async (e) => {{
                e.preventDefault();
                const btn = document.getElementById('submitDocxUploadBtn');
                const resultBox = document.getElementById('resultBox');
                const file = document.getElementById('docxFile').files[0];
                if (!file) return;

                btn.disabled = true;
                btn.textContent = 'Uploading...';
                resultBox.style.display = 'none';

                try {{
                    const formData = new FormData();
                    formData.append('docx_file', file);
                    const response = await fetch('/api/format-docx-upload', {{
                        method: 'POST',
                        body: formData
                    }});
                    const data = await response.json();

                    resultBox.style.display = 'block';
                    if (data.success) {{
                        resultBox.className = 'result success';
                        resultBox.innerHTML = `<strong>Success!</strong> Document formatted. <br><a href="${{data.file_url}}" target="_blank">Download ${{data.filename || 'File'}}</a>`;
                        setTimeout(() => window.location.reload(), 2000);
                    }} else {{
                        resultBox.className = 'result error';
                        resultBox.innerHTML = `<strong>Error!</strong><br><pre>${{data.message}}</pre>`;
                    }}
                }} catch (err) {{
                    resultBox.style.display = 'block';
                    resultBox.className = 'result error';
                    resultBox.textContent = 'Network error occurred.';
                }} finally {{
                    btn.disabled = false;
                    btn.textContent = 'Upload DOCX & Format';
                }}
            }});
            document.getElementById('pdfForm').addEventListener('submit', async (e) => {{
                e.preventDefault();
                const btn = document.getElementById('submitPdfBtn');
                const resultBox = document.getElementById('resultBox');
                
                const fileInput = document.getElementById('pdfFileInput');
                let source = document.getElementById('pdfSource').value;
                const instructions = document.getElementById('pdfInstructions').value;
                const theme = document.getElementById('pdfLayoutTheme').value;
                const iconography = document.getElementById('pdfIconography').value;
                const rules = document.getElementById('pdfContentRules').value;
                const format = document.getElementById('pdfTargetFormat').value;
                
                if (fileInput.files.length > 0) {{
                    source = await fileToBase64(fileInput.files[0]);
                }}
                
                if (!source) {{
                    alert("Please provide a PDF file or URL.");
                    return;
                }}
                
                btn.disabled = true;
                btn.textContent = 'Processing...';
                resultBox.style.display = 'none';
                
                try {{
                    const response = await fetch('/api/process-pdf', {{
                        method: 'POST',
                        headers: {{ 'Content-Type': 'application/json' }},
                        body: JSON.stringify({{ 
                            pdf_source: source, 
                            is_url: true,
                            instructions: instructions,
                            layout_theme: theme,
                            visual_iconography: iconography,
                            slide_content_rules: rules,
                            target_format: format,
                            api_key: document.getElementById('pdfApiKey') ? document.getElementById('pdfApiKey').value : ''
                        }})
                    }});
                    
                    const data = await response.json();
                    
                    resultBox.style.display = 'block';
                    if (data.success) {{
                        resultBox.className = 'result success';
                        resultBox.innerHTML = `<strong>Success!</strong> File generated. <br><a href="${{data.file_url}}" target="_blank">Download ${{data.filename || 'File'}}</a>`;
                        setTimeout(() => window.location.reload(), 2000);
                    }} else {{
                        resultBox.className = 'result error';
                        resultBox.innerHTML = `<strong>Error!</strong><br><pre>${{data.message}}</pre>`;
                    }}
                }} catch (err) {{
                    resultBox.style.display = 'block';
                    resultBox.className = 'result error';
                    resultBox.textContent = 'Network error occurred.';
                }} finally {{
                    btn.disabled = false;
                    btn.textContent = 'Process PDF';
                }}
            }});
            document.getElementById('pdfUploadForm').addEventListener('submit', async (e) => {{
                e.preventDefault();
                const btn = document.getElementById('submitPdfUploadBtn');
                const resultBox = document.getElementById('resultBox');
                const file = document.getElementById('pdfFile').files[0];

                if (!file) return;

                btn.disabled = true;
                btn.textContent = 'Uploading...';
                resultBox.style.display = 'none';

                try {{
                    const formData = new FormData();
                    formData.append('pdf_file', file);
                    formData.append('instructions', document.getElementById('pdfUploadInstructions').value || '');
                    formData.append('layout_theme', document.getElementById('pdfUploadLayoutTheme').value || '');
                    formData.append('visual_iconography', document.getElementById('pdfUploadIconography').value || '');
                    formData.append('slide_content_rules', document.getElementById('pdfUploadContentRules').value || '');
                    formData.append('target_format', document.getElementById('pdfUploadTargetFormat').value || 'pptx');
                    formData.append('api_key', document.getElementById('pdfUploadApiKey').value || '');

                    const response = await fetch('/api/process-pdf-upload', {{
                        method: 'POST',
                        body: formData
                    }});
                    const data = await response.json();

                    resultBox.style.display = 'block';
                    if (data.success) {{
                        resultBox.className = 'result success';
                        resultBox.innerHTML = `<strong>Success!</strong> File generated. <br><a href="${{data.file_url}}" target="_blank">Download ${{data.filename || 'File'}}</a>`;
                        setTimeout(() => window.location.reload(), 2000);
                    }} else {{
                        resultBox.className = 'result error';
                        resultBox.innerHTML = `<strong>Error!</strong><br><pre>${{data.message}}</pre>`;
                    }}
                }} catch (err) {{
                    resultBox.style.display = 'block';
                    resultBox.className = 'result error';
                    resultBox.textContent = 'Network error occurred.';
                }} finally {{
                    btn.disabled = false;
                    btn.textContent = 'Upload PDF & Process';
                }}
            }});
        </script>
    </body>
    </html>
    """
    return HTMLResponse(content=html_content)

@app.get("/api/stats")
async def get_stats():
    return stats

@app.get("/api/history")
async def get_history():
    return generation_history

@app.post("/api/generate")
async def api_generate(request: GenerateRequest):
    return generate_presentation(request.python_code, request.webhook_url)

@app.post("/api/image-to-pptx")
async def api_image_to_pptx(request: ImageRequest):
    return image_to_presentation(request.image_source, request.is_url, request.webhook_url)

@app.post("/api/format-docx")
async def api_format_docx(request: DocxRequest):
    return format_document(request.doc_source, request.is_url, request.webhook_url)

@app.post("/api/format-docx-upload")
async def api_format_docx_upload(docx_file: UploadFile = File(...), webhook_url: str = Form(None)):
    if not docx_file.filename.lower().endswith(".docx"):
        raise HTTPException(status_code=400, detail="Only .docx files are supported.")
    if docx_file.size and docx_file.size > MAX_UPLOAD_SIZE_BYTES:
        raise HTTPException(status_code=413, detail="File too large for this deployment.")
    temp_path = _persist_upload_to_tempfile(docx_file, ".docx")
    try:
        return format_document(temp_path, is_url=True, webhook_url=webhook_url)
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/api/format-docx-upload")
async def api_format_docx_upload(docx_file: UploadFile = File(...)):
    if not docx_file.filename.lower().endswith(".docx"):
        raise HTTPException(status_code=400, detail="Only .docx files are supported.")
    if docx_file.size and docx_file.size > MAX_UPLOAD_SIZE_BYTES:
        raise HTTPException(status_code=413, detail="File too large for this deployment.")
    temp_path = _persist_upload_to_tempfile(docx_file, ".docx")
    try:
        return format_document(temp_path, is_url=True)
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/api/process-pdf")
async def api_process_pdf(request: PdfRequest):
    return process_pdf_to_artifacts(
        request.pdf_source, 
        request.is_url, 
        request.instructions, 
        request.layout_theme, 
        request.visual_iconography, 
        request.slide_content_rules, 
        request.target_format,
        request.webhook_url,
        request.api_key
    )

@app.post("/api/process-pdf-upload")
async def api_process_pdf_upload(
    pdf_file: UploadFile = File(...),
    instructions: str = Form(""),
    layout_theme: str = Form(""),
    visual_iconography: str = Form(""),
    slide_content_rules: str = Form(""),
    target_format: str = Form("pptx"),
    webhook_url: str = Form(None),
    api_key: str = Form("")
):
    if not pdf_file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only .pdf files are supported.")
    if pdf_file.size and pdf_file.size > MAX_UPLOAD_SIZE_BYTES:
        raise HTTPException(status_code=413, detail="File too large for this deployment.")
    if target_format.lower() not in {"pptx", "docx"}:
        raise HTTPException(status_code=400, detail="target_format must be pptx or docx.")
    temp_path = _persist_upload_to_tempfile(pdf_file, ".pdf")
    try:
        return process_pdf_to_artifacts(
            temp_path,
            True,
            instructions,
            layout_theme,
            visual_iconography,
            slide_content_rules,
            target_format,
            webhook_url,
            api_key
        )
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.get("/downloads/{execution_id}/{filename}")
def download_file(execution_id: str, filename: str):
    file_path = os.path.join(OUTPUT_DIR, execution_id, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")
        
    with open(file_path, "rb") as f:
        data = f.read()
        
    media_type = "application/octet-stream"
    if filename.endswith(".pptx"):
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif filename.endswith(".docx"):
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        
    return Response(
        content=data,
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
